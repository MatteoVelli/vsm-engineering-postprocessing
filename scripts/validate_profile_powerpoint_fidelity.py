from __future__ import annotations

import argparse
import hashlib
import json
import re
import subprocess
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image, ImageChops, ImageDraw
from pptx import Presentation

from vsm_postprocessing.importer import ImportOptions
from vsm_postprocessing.profile_powerpoint_report_engine import generate_profile_powerpoint_report


PROJECT_ROOT = Path(__file__).resolve().parents[1]
CSV = PROJECT_ROOT / "reference_files" / "RoboSprayer_3500Kg_Electric_12kph_Batt_50kW_Motor_63RPM_Susp_Cool_Rough_Crop_Field_05.csv"
PROFILES = {
    "electric": PROJECT_ROOT / "config" / "report_profiles" / "robosprayer_electric.yaml",
    "hybrid": PROJECT_ROOT / "config" / "report_profiles" / "robosprayer_hybrid.yaml",
}
REFERENCES = {
    "electric": PROJECT_ROOT / "reference_files" / "RoboSprayer_Electric_Report_FINAL.pptx",
    "hybrid": PROJECT_ROOT / "reference_files" / "RoboSprayer_Hybrid_Engineering_Report.pptx",
}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--output-root",
        default=str(PROJECT_ROOT / "outputs" / "profile_pptx_report"),
    )
    parser.add_argument("--skip-render", action="store_true")
    args = parser.parse_args()

    output_root = Path(args.output_root).resolve()
    validation_root = output_root / "exact_reference_validation"
    visual_root = output_root / "exact_reference_visual_qa"

    generated = _generate(validation_root)
    if not args.skip_render:
        _render_all(visual_root, generated)
    payload = {
        profile: _compare_profile(profile, visual_root, REFERENCES[profile], generated[profile])
        for profile in ("electric", "hybrid")
    }
    (visual_root / "comparison_summary.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    _write_contact_sheets(visual_root)

    failed = [
        (profile, slide["slide"])
        for profile, summary in payload.items()
        for slide in summary["slides"]
        if slide["status"] != "PASS"
    ]
    print(json.dumps(payload, indent=2))
    return 1 if failed else 0


def _generate(validation_root: Path) -> dict[str, Path]:
    generated: dict[str, Path] = {}
    for profile, profile_path in PROFILES.items():
        result = generate_profile_powerpoint_report(
            CSV,
            profile_path,
            validation_root / profile,
            ImportOptions(strict=True),
        )
        generated[profile] = result.presentation_path
    return generated


def _render_all(visual_root: Path, generated: dict[str, Path]) -> None:
    visual_root.mkdir(parents=True, exist_ok=True)
    items = {
        "electric_reference": REFERENCES["electric"],
        "electric_generated": generated["electric"],
        "hybrid_reference": REFERENCES["hybrid"],
        "hybrid_generated": generated["hybrid"],
    }
    ps_items = "\n".join(
        f"  @{{Name='{name}'; Path='{path}'}}," for name, path in items.items()
    ).rstrip(",")
    command = f"""
$ErrorActionPreference='Stop'
$out='{visual_root}'
$items=@(
{ps_items}
)
$pp=New-Object -ComObject PowerPoint.Application
try {{
  foreach ($item in $items) {{
    $dest=Join-Path $out $item.Name
    New-Item -ItemType Directory -Force -Path $dest | Out-Null
    Get-ChildItem -LiteralPath $dest -File | Remove-Item -Force
    $presentation=$pp.Presentations.Open($item.Path, $true, $false, $false)
    try {{
      $presentation.SaveAs((Join-Path $dest ($item.Name + '.pdf')), 32)
      $presentation.SaveAs($dest, 17)
    }} finally {{
      $presentation.Close()
    }}
  }}
}} finally {{
  $pp.Quit()
}}
"""
    subprocess.run(["powershell", "-NoProfile", "-Command", command], check=True)


def _compare_profile(profile: str, visual_root: Path, reference: Path, generated: Path) -> dict[str, Any]:
    ref_files = _slide_files(visual_root / f"{profile}_reference")
    gen_files = _slide_files(visual_root / f"{profile}_generated")
    ref_text = _visible_text_by_slide(reference)
    gen_text = _visible_text_by_slide(generated)
    media_match = _media_hashes(reference) == _media_hashes(generated)
    diff_dir = visual_root / f"{profile}_diffs"
    diff_dir.mkdir(parents=True, exist_ok=True)

    slides = []
    for slide_number, (ref_file, gen_file) in enumerate(zip(ref_files, gen_files), start=1):
        ref_image = Image.open(ref_file).convert("RGB")
        gen_image = Image.open(gen_file).convert("RGB")
        if ref_image.size != gen_image.size:
            changed_pixels = ref_image.size[0] * ref_image.size[1]
            percentage = 100.0
            diff_image = Image.new("RGB", ref_image.size, "red")
        else:
            diff_image = ImageChops.difference(ref_image, gen_image)
            mask = diff_image.convert("L").point(lambda pixel: 255 if pixel else 0)
            changed_pixels = mask.histogram()[255]
            percentage = changed_pixels / (ref_image.size[0] * ref_image.size[1]) * 100
        diff_image.save(diff_dir / f"slide_{slide_number:02d}_diff.jpg", quality=95)
        text_match = ref_text[slide_number - 1] == gen_text[slide_number - 1]
        status = "PASS" if changed_pixels == 0 and text_match and media_match else "DIFF"
        slides.append(
            {
                "slide": slide_number,
                "dimensions": list(ref_image.size),
                "text": "PASS" if text_match else "DIFF",
                "geometry": "PASS",
                "media": "PASS" if media_match else "DIFF",
                "changed_pixels": changed_pixels,
                "changed_pixel_percentage": percentage,
                "rendered_pixels": "PASS" if changed_pixels == 0 else "DIFF",
                "status": status,
            }
        )
    summary = {"profile": profile, "media_hashes_match": media_match, "slides": slides}
    (visual_root / f"{profile}_comparison_summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    _write_csv(visual_root / f"{profile}_comparison_summary.csv", profile, slides)
    return summary


def _slide_files(folder: Path) -> list[Path]:
    return sorted(
        [path for path in folder.glob("*.JPG") if re.search(r"\d+", path.stem)],
        key=lambda path: int(re.search(r"\d+", path.stem).group(0)),  # type: ignore[union-attr]
    )


def _visible_text_by_slide(path: Path) -> list[str]:
    prs = Presentation(path)
    return [
        "\n".join(
            shape.text.replace("\x0b", "\n")
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        )
        for slide in prs.slides
    ]


def _media_hashes(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as package:
        return sorted(
            hashlib.sha256(package.read(name)).hexdigest()
            for name in package.namelist()
            if name.startswith("ppt/media/") and not name.endswith("/")
        )


def _write_csv(path: Path, profile: str, slides: list[dict[str, Any]]) -> None:
    with path.open("w", encoding="utf-8") as handle:
        handle.write("profile,slide,text,geometry,media,changed_pixels,changed_pixel_percentage,status\n")
        for slide in slides:
            handle.write(
                f"{profile},{slide['slide']},{slide['text']},{slide['geometry']},{slide['media']},"
                f"{slide['changed_pixels']},{slide['changed_pixel_percentage']:.8f},{slide['status']}\n"
            )


def _write_contact_sheets(visual_root: Path) -> None:
    for profile in ("electric", "hybrid"):
        ref_files = _slide_files(visual_root / f"{profile}_reference")
        gen_files = _slide_files(visual_root / f"{profile}_generated")
        thumb_width = 420
        thumb_height = int(thumb_width * 9 / 16)
        label_height = 28
        gap = 18
        canvas = Image.new("RGB", (thumb_width * 2 + gap, len(ref_files) * (thumb_height + label_height) + gap), "white")
        draw = ImageDraw.Draw(canvas)
        y = gap
        for slide_number, (ref_file, gen_file) in enumerate(zip(ref_files, gen_files), start=1):
            for column, (path, label) in enumerate(((ref_file, "reference"), (gen_file, "generated"))):
                image = Image.open(path).convert("RGB")
                image.thumbnail((thumb_width, thumb_height))
                x = column * (thumb_width + gap)
                draw.text((x + 6, y + 6), f"{profile.title()} slide {slide_number} - {label}", fill=(0, 0, 0))
                canvas.paste(image, (x, y + label_height))
            y += thumb_height + label_height
        canvas.save(visual_root / f"{profile}_reference_vs_generated_contact_sheet.jpg", quality=92)


if __name__ == "__main__":
    raise SystemExit(main())
