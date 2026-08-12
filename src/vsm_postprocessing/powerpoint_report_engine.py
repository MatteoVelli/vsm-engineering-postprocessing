from __future__ import annotations

import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .errors import ConfigurationError, PowerPointReportError
from .importer import ImportOptions
from .plotting_engine import PlottingResult, render_plots
from .statistics_engine import StatisticResult, StatisticsResult, calculate_statistics
from .utils import client_display_filename, sha256_file
from .version import __version__


@dataclass(frozen=True)
class PowerPointSlideDefinition:
    slide_id: str
    slide_type: str
    title: str
    subtitle: str | None
    statistic_ids: tuple[str, ...]
    plot_ids: tuple[str, ...]
    footer: str | None
    body: tuple[str, ...] = ()
    layout: str | None = None
    appendix: bool = False


@dataclass(frozen=True)
class PowerPointTheme:
    title_font_size: float = 28.0
    subtitle_font_size: float = 14.0
    kpi_label_font_size: float = 10.0
    kpi_value_font_size: float = 19.0
    body_font_size: float = 11.0
    accent_fill: str = "FFF2CC"
    accent_border: str = "D6B656"
    text_color: str = "111111"
    muted_text_color: str = "666666"
    background_color: str = "FFFFFF"
    rule_color: str = "B8C2CC"


@dataclass(frozen=True)
class PowerPointReportConfig:
    version: int
    title: str
    subtitle: str | None
    output_filename: str
    footer: str | None
    author: str
    subject: str | None
    keywords: str | None
    comments: str | None
    theme: PowerPointTheme
    slides: tuple[PowerPointSlideDefinition, ...]


@dataclass
class PowerPointReportResult:
    presentation_path: Path
    manifest_path: Path
    summary_path: Path
    plot_assets_dir: Path
    config_path: Path
    config: PowerPointReportConfig
    statistics_result: StatisticsResult
    plotting_result: PlottingResult

    @property
    def sample_count(self) -> int:
        return self.statistics_result.sample_count

    @property
    def slide_count(self) -> int:
        return len(self.config.slides)

    @property
    def plot_count(self) -> int:
        return len({pid for slide in self.config.slides for pid in slide.plot_ids})

    @property
    def statistic_count(self) -> int:
        return len({sid for slide in self.config.slides for sid in slide.statistic_ids})

    @property
    def displayed_kpi_count(self) -> int:
        return sum(len(slide.statistic_ids) for slide in self.config.slides)

_ALLOWED_SLIDE_TYPES = {"cover", "overview", "summary", "kpi_grid", "plot_full", "plot_pair", "conclusion"}


def load_powerpoint_report_config(path: str | Path) -> PowerPointReportConfig:
    config_path = Path(path).expanduser().resolve()
    if not config_path.exists():
        raise ConfigurationError(f"PowerPoint-report configuration does not exist: {config_path}")
    if not config_path.is_file():
        raise ConfigurationError(f"PowerPoint-report configuration is not a file: {config_path}")
    try:
        raw = yaml.safe_load(config_path.read_text(encoding="utf-8"))
    except UnicodeDecodeError as exc:
        raise ConfigurationError(f"PowerPoint-report configuration must be UTF-8: {config_path}") from exc
    except yaml.YAMLError as exc:
        raise ConfigurationError(f"Invalid YAML in PowerPoint-report configuration '{config_path}': {exc}") from exc
    if not isinstance(raw, dict):
        raise ConfigurationError("PowerPoint-report configuration root must be a YAML mapping")
    _reject_unknown_keys(raw, {"version", "presentation", "theme", "slides"}, "root")
    if raw.get("version") != 1:
        raise ConfigurationError("PowerPoint-report configuration 'version' must be 1")

    presentation = raw.get("presentation", {})
    if not isinstance(presentation, dict):
        raise ConfigurationError("presentation must be a YAML mapping")
    _reject_unknown_keys(
        presentation,
        {"title", "subtitle", "output_filename", "footer", "author", "subject", "keywords", "comments"},
        "presentation",
    )
    title = _nonempty_string(presentation.get("title", "VSM Engineering Report"), "presentation.title")
    subtitle = _optional_string(presentation.get("subtitle"), "presentation.subtitle")
    footer = _optional_string(presentation.get("footer"), "presentation.footer")
    author = _nonempty_string(presentation.get("author", "VSM Engineering Post-Processing Tool"), "presentation.author")
    subject = _optional_string(presentation.get("subject"), "presentation.subject")
    keywords = _optional_string(presentation.get("keywords"), "presentation.keywords")
    comments = _optional_string(presentation.get("comments"), "presentation.comments")
    output_filename = _pptx_filename(presentation.get("output_filename", "vsm_engineering_report.pptx"))

    theme = _load_theme(raw.get("theme", {}))
    slides_raw = raw.get("slides")
    if not isinstance(slides_raw, list) or not slides_raw:
        raise ConfigurationError("slides must be a non-empty YAML list")
    slides = tuple(_load_slide(item, index) for index, item in enumerate(slides_raw, start=1))
    duplicate_ids = _duplicates(slide.slide_id for slide in slides)
    if duplicate_ids:
        raise ConfigurationError("slides contains duplicate slide IDs: " + ", ".join(duplicate_ids))
    return PowerPointReportConfig(
        version=1,
        title=title,
        subtitle=subtitle,
        output_filename=output_filename,
        footer=footer,
        author=author,
        subject=subject,
        keywords=keywords,
        comments=comments,
        theme=theme,
        slides=slides,
    )


def generate_powerpoint_report(
    input_file: str | Path,
    report_config_file: str | Path,
    statistics_config_file: str | Path,
    plotting_config_file: str | Path,
    output_dir: str | Path,
    import_options: ImportOptions | None = None,
    math_config_file: str | Path | None = None,
) -> PowerPointReportResult:
    destination = Path(output_dir).expanduser().resolve()
    destination.mkdir(parents=True, exist_ok=True)
    plot_assets_dir = destination / "plot_assets"
    plot_assets_dir.mkdir(parents=True, exist_ok=True)
    statistics_result = calculate_statistics(
        input_file,
        statistics_config_file,
        import_options,
        math_config_file=math_config_file,
    )
    plotting_result = render_plots(
        input_file,
        plotting_config_file,
        plot_assets_dir,
        import_options,
        math_config_file=math_config_file,
    )
    return build_powerpoint_report(
        statistics_result,
        plotting_result,
        report_config_file,
        destination,
        plot_assets_dir=plot_assets_dir,
    )


def build_powerpoint_report(
    statistics_result: StatisticsResult,
    plotting_result: PlottingResult,
    report_config_file: str | Path,
    output_dir: str | Path,
    *,
    plot_assets_dir: str | Path,
) -> PowerPointReportResult:
    config_path = Path(report_config_file).expanduser().resolve()
    config = load_powerpoint_report_config(config_path)
    destination = Path(output_dir).expanduser().resolve()
    destination.mkdir(parents=True, exist_ok=True)
    assets = Path(plot_assets_dir).expanduser().resolve()

    statistics_by_id = {item.statistic_id: item for item in statistics_result.statistics}
    plots_by_id = {item.plot_id: item for item in plotting_result.rendered_plots}
    required_statistics = {sid for slide in config.slides for sid in slide.statistic_ids}
    required_plots = {pid for slide in config.slides for pid in slide.plot_ids}
    missing_statistics = sorted(required_statistics - set(statistics_by_id))
    missing_plots = sorted(required_plots - set(plots_by_id))
    if missing_statistics:
        raise PowerPointReportError("Configured PowerPoint statistic IDs were not found: " + ", ".join(missing_statistics))
    if missing_plots:
        raise PowerPointReportError("Configured PowerPoint plot IDs were not found: " + ", ".join(missing_plots))

    for plot_id in required_plots:
        path = _plot_path(plots_by_id[plot_id].output_file, assets)
        if not path.exists():
            raise PowerPointReportError(f"Plot asset for '{plot_id}' does not exist: {path}")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = config.title
    prs.core_properties.subject = config.subject or config.subtitle or "Deterministic VSM engineering report"
    prs.core_properties.author = config.author
    if config.keywords:
        prs.core_properties.keywords = config.keywords
    prs.core_properties.comments = config.comments or f"Generated by VSM Engineering Post-Processing Tool v{__version__}"

    for slide_number, slide_def in enumerate(config.slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _apply_slide_background(slide, config.theme)
        if slide_def.slide_type == "cover":
            _render_cover_slide(slide, slide_def, statistics_by_id, statistics_result, config, slide_number)
        elif slide_def.slide_type == "overview":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_overview_slide(slide, slide_def, statistics_by_id, config.theme)
        elif slide_def.slide_type == "summary":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_summary_slide(slide, slide_def, statistics_by_id, statistics_result, config.theme)
        elif slide_def.slide_type == "kpi_grid":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_kpi_grid_slide(slide, slide_def, statistics_by_id, config.theme)
        elif slide_def.slide_type == "plot_full":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_plot_full_slide(slide, slide_def, statistics_by_id, plots_by_id, assets, config.theme)
        elif slide_def.slide_type == "plot_pair":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_plot_pair_slide(slide, slide_def, statistics_by_id, plots_by_id, assets, config.theme)
        elif slide_def.slide_type == "conclusion":
            _add_slide_title(slide, slide_def.title, slide_def.subtitle, config.theme)
            _render_conclusion_slide(slide, slide_def, statistics_by_id, config.theme)
        else:  # defensive: loader already validates
            raise PowerPointReportError(f"Unsupported PowerPoint slide type: {slide_def.slide_type}")
        _add_footer(slide, slide_def.footer or config.footer, slide_number, len(config.slides), config.theme)

    presentation_path = destination / config.output_filename
    try:
        prs.save(presentation_path)
    except Exception as exc:  # python-pptx exposes several low-level exceptions
        raise PowerPointReportError(f"Could not save PowerPoint report '{presentation_path}': {exc}") from exc

    manifest_path = destination / "powerpoint_report_manifest.json"
    summary_path = destination / "powerpoint_report_summary.txt"
    result = PowerPointReportResult(
        presentation_path=presentation_path,
        manifest_path=manifest_path,
        summary_path=summary_path,
        plot_assets_dir=assets,
        config_path=config_path,
        config=config,
        statistics_result=statistics_result,
        plotting_result=plotting_result,
    )
    _write_metadata(result)
    return result


def _render_summary_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    statistics_result: StatisticsResult,
    theme: PowerPointTheme,
) -> None:
    source_name = client_display_filename(statistics_result.dataset.source_path)
    source_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.42))
    p = source_box.text_frame.paragraphs[0]
    p.text = f"Source: {source_name}  |  Samples: {statistics_result.sample_count}"
    p.font.size = Pt(theme.body_font_size)
    p.font.color.rgb = _rgb(theme.muted_text_color)
    p.alignment = PP_ALIGN.CENTER
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    _add_kpi_grid(slide, stats, y=2.15, theme=theme, max_columns=3, card_height=1.45)


def _render_cover_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    statistics_result: StatisticsResult,
    config: PowerPointReportConfig,
    slide_number: int,
) -> None:
    theme = config.theme
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.65), Inches(11.85), Inches(0.75))
    p = title_box.text_frame.paragraphs[0]
    p.text = definition.title
    p.font.size = Pt(30)
    p.font.color.rgb = _rgb(theme.text_color)
    p.font.bold = False
    p.alignment = PP_ALIGN.CENTER
    if definition.subtitle:
        subtitle = slide.shapes.add_textbox(Inches(0.95), Inches(1.45), Inches(11.45), Inches(0.45))
        sp = subtitle.text_frame.paragraphs[0]
        sp.text = definition.subtitle
        sp.font.size = Pt(theme.subtitle_font_size + 2)
        sp.font.color.rgb = _rgb(theme.muted_text_color)
        sp.alignment = PP_ALIGN.CENTER
    _add_rule(slide, 2.08, theme, x=3.15, w=7.05)
    if definition.body:
        _add_bullet_panel(
            slide,
            definition.body,
            1.15,
            2.45,
            5.25,
            2.85,
            theme,
            title="Mission Definition",
            body_font_size=11,
        )
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    if stats:
        _add_kpi_grid(slide, stats[:6], y=2.55, theme=theme, max_columns=2, card_height=0.86, x=7.0, total_w=5.35)
    meta = [
        f"Samples: {statistics_result.sample_count:,}",
        f"Tool version: v{__version__}",
    ]
    meta_box = slide.shapes.add_textbox(Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.28))
    mp = meta_box.text_frame.paragraphs[0]
    mp.text = "  |  ".join(meta)
    mp.font.size = Pt(11)
    mp.font.color.rgb = _rgb(theme.muted_text_color)
    mp.alignment = PP_ALIGN.CENTER


def _render_overview_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    theme: PowerPointTheme,
) -> None:
    left = definition.body[:5]
    right = definition.body[5:]
    _add_bullet_panel(slide, left, 0.72, 1.65, 5.75, 3.95, theme, title="Mission Structure")
    _add_bullet_panel(slide, right, 6.85, 1.65, 5.75, 3.95, theme, title="Hybrid System Context")
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    if stats:
        _add_kpi_strip(slide, stats, y=5.92, theme=theme)


def _render_kpi_grid_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    theme: PowerPointTheme,
) -> None:
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    _add_kpi_grid(slide, stats, y=1.62, theme=theme, max_columns=4, card_height=1.02)
    if definition.body:
        _add_body_text(slide, definition.body, 0.9, 6.25, 11.5, 0.48, theme)


def _render_plot_full_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    plots_by_id: dict[str, Any],
    assets: Path,
    theme: PowerPointTheme,
) -> None:
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    if stats:
        _add_kpi_strip(slide, stats, y=1.32, theme=theme)
    path = _plot_path(plots_by_id[definition.plot_ids[0]].output_file, assets)
    _add_contained_picture(slide, path, 1.05, 2.18, 11.25, 4.55)
    if definition.body:
        _add_body_text(slide, definition.body, 0.9, 6.78, 11.4, 0.28, theme)


def _render_plot_pair_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    plots_by_id: dict[str, Any],
    assets: Path,
    theme: PowerPointTheme,
) -> None:
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    if stats:
        _add_kpi_strip(slide, stats, y=1.45, theme=theme)
    plot_y = 2.2 if stats else 1.62
    plot_h = 4.25 if stats else 4.95
    plot_ids = list(definition.plot_ids)
    if len(plot_ids) == 1:
        path = _plot_path(plots_by_id[plot_ids[0]].output_file, assets)
        _add_contained_picture(slide, path, 1.35, plot_y, 10.65, plot_h)
    elif len(plot_ids) == 2:
        positions = [(0.42, plot_y, 6.15, plot_h), (6.76, plot_y, 6.15, plot_h)]
        for plot_id, (x, y, w, h) in zip(plot_ids, positions):
            path = _plot_path(plots_by_id[plot_id].output_file, assets)
            _add_contained_picture(slide, path, x, y, w, h)
    if definition.body:
        _add_body_text(slide, definition.body, 0.75, 6.72, 11.85, 0.28, theme)


def _render_conclusion_slide(
    slide: Any,
    definition: PowerPointSlideDefinition,
    statistics_by_id: dict[str, StatisticResult],
    theme: PowerPointTheme,
) -> None:
    stats = [statistics_by_id[sid] for sid in definition.statistic_ids]
    _add_kpi_grid(slide, stats[:6], y=1.45, theme=theme, max_columns=3, card_height=1.02)
    if definition.body:
        _add_bullet_panel(slide, definition.body, 1.05, 4.45, 11.25, 1.75, theme, title="Deterministic Summary")


def _add_slide_title(slide: Any, title: str, subtitle: str | None, theme: PowerPointTheme) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.26), Inches(12.63), Inches(0.72))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(theme.title_font_size)
    p.font.bold = False
    p.font.color.rgb = _rgb(theme.text_color)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.93), Inches(12.33), Inches(0.34))
        sp = subtitle_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(theme.subtitle_font_size)
        sp.font.color.rgb = _rgb(theme.muted_text_color)
    _add_rule(slide, 1.29, theme, x=4.55, w=4.23)


def _add_kpi_strip(slide: Any, statistics: list[StatisticResult], *, y: float, theme: PowerPointTheme) -> None:
    if not statistics:
        return
    max_items = min(len(statistics), 6)
    stats = statistics[:max_items]
    total_w = 12.1
    gap = 0.08
    card_w = (total_w - gap * (len(stats) - 1)) / len(stats)
    x = 0.62
    for stat in stats:
        _add_kpi_card(slide, stat, x, y, card_w, 0.78, theme, compact=True)
        x += card_w + gap


def _add_kpi_grid(
    slide: Any,
    statistics: list[StatisticResult],
    *,
    y: float,
    theme: PowerPointTheme,
    max_columns: int,
    card_height: float,
    x: float | None = None,
    total_w: float = 11.5,
) -> None:
    if not statistics:
        return
    stats = statistics[:6]
    columns = min(max_columns, len(stats))
    gap_x = 0.25
    gap_y = 0.24
    card_w = (total_w - gap_x * (columns - 1)) / columns
    rows = math.ceil(len(stats) / columns)
    start_x = x if x is not None else (13.333333 - total_w) / 2
    for index, stat in enumerate(stats):
        row, col = divmod(index, columns)
        x = start_x + col * (card_w + gap_x)
        card_y = y + row * (card_height + gap_y)
        _add_kpi_card(slide, stat, x, card_y, card_w, card_height, theme, compact=False)


def _add_kpi_card(
    slide: Any,
    stat: StatisticResult,
    x: float,
    y: float,
    w: float,
    h: float,
    theme: PowerPointTheme,
    *,
    compact: bool,
) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme.accent_fill)
    shape.line.color.rgb = _rgb(theme.accent_border)
    shape.line.width = Pt(0.8)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.text = _format_statistic_label(stat)
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(theme.kpi_label_font_size if compact else theme.kpi_label_font_size + 1)
    p1.font.color.rgb = _rgb(theme.text_color)
    p2 = tf.add_paragraph()
    p2.text = _format_statistic(stat)
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(theme.kpi_value_font_size if compact else theme.kpi_value_font_size + 3)
    p2.font.bold = True
    p2.font.color.rgb = _rgb(theme.text_color)


def _apply_slide_background(slide: Any, theme: PowerPointTheme) -> None:
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb(theme.background_color)


def _add_rule(slide: Any, y: float, theme: PowerPointTheme, *, x: float, w: float) -> None:
    accent = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.035))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.accent_border)
    accent.line.fill.background()


def _add_bullet_panel(
    slide: Any,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    theme: PowerPointTheme,
    *,
    title: str | None = None,
    body_font_size: float | None = None,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(theme.body_font_size + 2)
        p.font.bold = True
        p.font.color.rgb = _rgb(theme.text_color)
        first = False
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(body_font_size if body_font_size is not None else theme.body_font_size + 1)
        p.font.color.rgb = _rgb(theme.text_color)
        p.space_after = Pt(5)


def _add_body_text(
    slide: Any,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    theme: PowerPointTheme,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    p = box.text_frame.paragraphs[0]
    p.text = "  |  ".join(items)
    p.font.size = Pt(8.5)
    p.font.color.rgb = _rgb(theme.muted_text_color)
    p.alignment = PP_ALIGN.CENTER


def _add_footer(slide: Any, footer: str | None, slide_number: int, slide_count: int, theme: PowerPointTheme) -> None:
    if footer:
        box = slide.shapes.add_textbox(Inches(0.38), Inches(7.12), Inches(11.8), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = f"{footer} | v{__version__}"
        p.font.size = Pt(8)
        p.font.color.rgb = _rgb(theme.muted_text_color)
    num = slide.shapes.add_textbox(Inches(12.25), Inches(7.08), Inches(0.65), Inches(0.24))
    p = num.text_frame.paragraphs[0]
    p.text = f"{slide_number}/{slide_count}"
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(8)
    p.font.color.rgb = _rgb(theme.muted_text_color)


def _add_contained_picture(slide: Any, image_path: Path, x: float, y: float, w: float, h: float) -> None:
    try:
        with Image.open(image_path) as image:
            width_px, height_px = image.size
    except Exception as exc:
        raise PowerPointReportError(f"Could not read plot image '{image_path}': {exc}") from exc
    if width_px <= 0 or height_px <= 0:
        raise PowerPointReportError(f"Plot image has invalid dimensions: {image_path}")
    image_ratio = width_px / height_px
    box_ratio = w / h
    if image_ratio >= box_ratio:
        final_w = w
        final_h = w / image_ratio
        final_x = x
        final_y = y + (h - final_h) / 2
    else:
        final_h = h
        final_w = h * image_ratio
        final_x = x + (w - final_w) / 2
        final_y = y
    slide.shapes.add_picture(
        str(image_path),
        Inches(final_x),
        Inches(final_y),
        width=Inches(final_w),
        height=Inches(final_h),
    )


def _plot_path(output_file: str, assets: Path) -> Path:
    candidate = Path(output_file)
    if candidate.is_absolute():
        return candidate
    return assets / candidate.name


def _write_metadata(result: PowerPointReportResult) -> None:
    used_stats = sorted({sid for slide in result.config.slides for sid in slide.statistic_ids})
    used_plots = sorted({pid for slide in result.config.slides for pid in slide.plot_ids})
    plot_files = {
        item.plot_id: item.output_file
        for item in result.plotting_result.rendered_plots
        if item.plot_id in used_plots
    }
    manifest = {
        "status": "PASS",
        "source_file": str(result.statistics_result.dataset.source_path),
        "source_sha256": sha256_file(result.statistics_result.dataset.source_path),
        "configuration_file": str(result.config_path),
        "configuration_sha256": sha256_file(result.config_path),
        "presentation": str(result.presentation_path),
        "sample_count": result.sample_count,
        "slide_count": result.slide_count,
        "statistic_count": result.statistic_count,
        "displayed_kpi_count": result.displayed_kpi_count,
        "plot_count": result.plot_count,
        "matplotlib_image_count": len(used_plots),
        "matplotlib_image_placement_count": sum(len(slide.plot_ids) for slide in result.config.slides),
        "appendix_slide_count": sum(1 for slide in result.config.slides if slide.appendix),
        "document_properties": {
            "title": result.config.title,
            "subject": result.config.subject or result.config.subtitle or "Deterministic VSM engineering report",
            "author": result.config.author,
            "keywords": result.config.keywords or "",
            "comments": result.config.comments or f"Generated by VSM Engineering Post-Processing Tool v{__version__}",
        },
        "footer_format": f"{result.config.footer or ''} | v{__version__} | slide/total",
        "statistics_used": used_stats,
        "plots_used": used_plots,
        "plot_files": plot_files,
        "slides": [
            {
                "slide_id": slide.slide_id,
                "slide_type": slide.slide_type,
                "title": slide.title,
                "appendix": slide.appendix,
                "statistics": list(slide.statistic_ids),
                "plots": list(slide.plot_ids),
            }
            for slide in result.config.slides
        ],
    }
    result.manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    lines = [
        "VSM POWERPOINT REPORT",
        "=====================",
        "Status: PASS",
        f"Source: {result.statistics_result.dataset.source_path}",
        f"Samples: {result.sample_count}",
        f"Slides: {result.slide_count}",
        f"Statistics used: {result.statistic_count}",
        f"Displayed KPIs: {result.displayed_kpi_count}",
        f"Plots used: {result.plot_count}",
        f"Appendix slides: {sum(1 for slide in result.config.slides if slide.appendix)}",
        f"Presentation: {result.presentation_path}",
        "",
        "Slides:",
    ]
    for index, slide in enumerate(result.config.slides, start=1):
        lines.append(
            f"{index:02d}. {slide.slide_id} | {slide.slide_type} | "
            f"statistics={len(slide.statistic_ids)} | plots={len(slide.plot_ids)}"
        )
    result.summary_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _load_theme(raw: object) -> PowerPointTheme:
    if raw is None:
        raw = {}
    if not isinstance(raw, dict):
        raise ConfigurationError("theme must be a YAML mapping")
    allowed = {
        "title_font_size", "subtitle_font_size", "kpi_label_font_size", "kpi_value_font_size",
        "body_font_size", "accent_fill", "accent_border", "text_color", "muted_text_color",
        "background_color", "rule_color",
    }
    _reject_unknown_keys(raw, allowed, "theme")
    defaults = PowerPointTheme()
    numeric_fields = {
        "title_font_size": defaults.title_font_size,
        "subtitle_font_size": defaults.subtitle_font_size,
        "kpi_label_font_size": defaults.kpi_label_font_size,
        "kpi_value_font_size": defaults.kpi_value_font_size,
        "body_font_size": defaults.body_font_size,
    }
    values: dict[str, Any] = {}
    for key, default in numeric_fields.items():
        value = raw.get(key, default)
        if isinstance(value, bool) or not isinstance(value, (int, float)) or not 6 <= float(value) <= 60:
            raise ConfigurationError(f"theme.{key} must be a number from 6 to 60")
        values[key] = float(value)
    for key, default in {
        "accent_fill": defaults.accent_fill,
        "accent_border": defaults.accent_border,
        "text_color": defaults.text_color,
        "muted_text_color": defaults.muted_text_color,
        "background_color": defaults.background_color,
        "rule_color": defaults.rule_color,
    }.items():
        values[key] = _hex_color(raw.get(key, default), f"theme.{key}")
    return PowerPointTheme(**values)


def _load_slide(raw: object, index: int) -> PowerPointSlideDefinition:
    context = f"slides[{index}]"
    if not isinstance(raw, dict):
        raise ConfigurationError(f"{context} must be a YAML mapping")
    _reject_unknown_keys(
        raw,
        {"slide_id", "type", "title", "subtitle", "statistics", "plots", "footer", "body", "layout", "appendix"},
        context,
    )
    slide_id = _identifier(raw.get("slide_id"), f"{context}.slide_id")
    slide_type = _nonempty_string(raw.get("type"), f"{context}.type")
    if slide_type not in _ALLOWED_SLIDE_TYPES:
        raise ConfigurationError(f"{context}.type must be one of: " + ", ".join(sorted(_ALLOWED_SLIDE_TYPES)))
    title = _nonempty_string(raw.get("title"), f"{context}.title")
    subtitle = _optional_string(raw.get("subtitle"), f"{context}.subtitle")
    footer = _optional_string(raw.get("footer"), f"{context}.footer")
    body = tuple(_string_list(raw.get("body", []), f"{context}.body", allow_empty=True))
    layout = _optional_string(raw.get("layout"), f"{context}.layout")
    appendix = raw.get("appendix", False)
    if not isinstance(appendix, bool):
        raise ConfigurationError(f"{context}.appendix must be true or false")
    statistic_ids = tuple(_string_list(raw.get("statistics", []), f"{context}.statistics", allow_empty=True))
    plot_ids = tuple(_string_list(raw.get("plots", []), f"{context}.plots", allow_empty=True))
    if slide_type in {"cover", "overview", "conclusion"}:
        if plot_ids:
            raise ConfigurationError(f"{context}.plots must be empty for {slide_type} slides")
        if len(statistic_ids) > 8:
            raise ConfigurationError(f"{context}.statistics supports at most 8 KPI statistics")
    if slide_type == "kpi_grid":
        if plot_ids:
            raise ConfigurationError(f"{context}.plots must be empty for kpi_grid slides")
        if not 1 <= len(statistic_ids) <= 12:
            raise ConfigurationError(f"{context}.statistics must contain one to twelve KPI statistics")
    if slide_type == "summary":
        if plot_ids:
            raise ConfigurationError(f"{context}.plots must be empty for summary slides")
        if len(statistic_ids) > 6:
            raise ConfigurationError(f"{context}.statistics supports at most 6 KPI statistics")
    if slide_type == "plot_full":
        if len(plot_ids) != 1:
            raise ConfigurationError(f"{context}.plots must contain exactly one plot ID for plot_full slides")
        if len(statistic_ids) > 6:
            raise ConfigurationError(f"{context}.statistics supports at most 6 KPI statistics")
    if slide_type == "plot_pair":
        if not 1 <= len(plot_ids) <= 2:
            raise ConfigurationError(f"{context}.plots must contain one or two plot IDs for plot_pair slides")
        if len(statistic_ids) > 6:
            raise ConfigurationError(f"{context}.statistics supports at most 6 KPI statistics")
    return PowerPointSlideDefinition(
        slide_id=slide_id,
        slide_type=slide_type,
        title=title,
        subtitle=subtitle,
        statistic_ids=statistic_ids,
        plot_ids=plot_ids,
        footer=footer,
        body=body,
        layout=layout,
        appendix=appendix,
    )


def _format_statistic(stat: StatisticResult) -> str:
    value = stat.value
    if not math.isfinite(value):
        value_text = "n/a"
    else:
        unit = _display_unit(stat.channel_unit)
        magnitude = abs(value)
        if unit == "rpm":
            value_text = f"{value:,.0f}"
        elif magnitude >= 10000:
            value_text = f"{value:,.0f}"
        elif magnitude >= 1000:
            value_text = f"{value:,.1f}"
        elif magnitude >= 100:
            value_text = f"{value:.2f}"
        elif unit in {"km", "min", "kg", "kW", "kWh", "%", "kph", "Nm", "l/h"}:
            value_text = f"{value:.2f}"
        else:
            value_text = f"{value:.3f}"
        if unit not in {"km", "min", "kg", "kW", "kWh", "%", "kph", "Nm", "l/h"}:
            value_text = value_text.rstrip("0").rstrip(".") if "." in value_text else value_text
    unit = _display_unit(stat.channel_unit)
    return f"{value_text} {unit}" if unit else value_text


def _format_statistic_label(stat: StatisticResult) -> str:
    text = stat.display_name
    return text.split("[", 1)[0].strip() if "[" in text else text


def _display_unit(unit: str | None) -> str | None:
    if not unit:
        return None
    normalized = unit.strip()
    return {"Km": "km", "Kg": "kg", "KW": "kW", "KWh": "kWh"}.get(normalized, normalized)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _hex_color(value: object, context: str) -> str:
    if isinstance(value, int) and not isinstance(value, bool):
        text = str(value)
    elif isinstance(value, str):
        text = value.strip().lstrip("#").upper()
    else:
        raise ConfigurationError(f"{context} must be a six-digit hexadecimal color")
    if len(text) != 6 or any(ch not in "0123456789ABCDEF" for ch in text):
        raise ConfigurationError(f"{context} must be a six-digit hexadecimal color")
    return text


def _pptx_filename(value: object) -> str:
    text = _nonempty_string(value, "presentation.output_filename")
    if Path(text).name != text or not text.lower().endswith(".pptx"):
        raise ConfigurationError("presentation.output_filename must be a plain .pptx filename")
    return text


def _nonempty_string(value: object, context: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise ConfigurationError(f"{context} must be a non-empty string")
    return value.strip()


def _optional_string(value: object, context: str) -> str | None:
    if value is None:
        return None
    return _nonempty_string(value, context)


def _identifier(value: object, context: str) -> str:
    text = _nonempty_string(value, context)
    if not all(ch.isalnum() or ch in {"_", "-"} for ch in text):
        raise ConfigurationError(f"{context} may contain only letters, numbers, '_' and '-'")
    return text


def _string_list(value: object, context: str, *, allow_empty: bool) -> list[str]:
    if not isinstance(value, list):
        raise ConfigurationError(f"{context} must be a YAML list")
    if not allow_empty and not value:
        raise ConfigurationError(f"{context} must not be empty")
    result: list[str] = []
    for item in value:
        if not isinstance(item, str) or not item.strip():
            raise ConfigurationError(f"{context} must contain only non-empty strings")
        text = item.strip()
        if text in result:
            raise ConfigurationError(f"{context} contains duplicate values")
        result.append(text)
    return result


def _duplicates(values: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    duplicates: list[str] = []
    for value in values:
        if value in seen and value not in duplicates:
            duplicates.append(value)
        seen.add(value)
    return duplicates


def _reject_unknown_keys(mapping: dict[str, Any], allowed: set[str], context: str) -> None:
    unknown = sorted(set(mapping) - allowed)
    if unknown:
        raise ConfigurationError(f"Unknown key(s) in {context}: " + ", ".join(unknown))
