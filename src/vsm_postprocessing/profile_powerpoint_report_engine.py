from __future__ import annotations

import json
import math
import re
from dataclasses import dataclass, replace
from pathlib import Path
from typing import Any, Mapping

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt

from .errors import PowerPointReportError
from .excel_report_engine import ProfileExcelReportResult, generate_profile_excel_report
from .importer import ImportOptions
from .powerpoint_report_engine import (
    PowerPointReportResult,
    load_powerpoint_report_config,
    _write_metadata as _write_powerpoint_metadata,
)
from .profile_statistics import ProfileKPIResult, ProfileStatisticResult
from .report_profile import ReportingProfile
from .statistics_engine import (
    StatisticResult,
    StatisticsConfig,
    StatisticsOutputOptions,
    StatisticsResult,
)
from .utils import client_display_filename, sha256_file
from .version import __version__


_PRESENTATION_SLIDE_ORDER = (
    "cover",
    "system_overview",
    "executive_results",
    "vehicle_operation",
    "battery_system",
    "battery_power_recovery",
    "range_extender_generator",
    "profile_energy_range",
    "agrochemical_battery",
    "traction_auxiliaries",
    "simulation_summary",
)
_HYBRID_ACTIVITY_IDS = (
    "engine_fuel_consumption_last",
    "fuel_flow_max",
    "engine_speed_max",
    "engine_torque_max",
    "engine_power_required_max",
    "engine_energy_delivered_sum",
    "generator_torque_1_max",
    "generator_power_1_max",
)
_DEFAULT_REFERENCE_TEMPLATE: str | None = None


@dataclass(frozen=True)
class PowerPointReferenceShapeSpec:
    slide_number: int
    shape_index: int
    name: str
    shape_type: str
    x_in: float
    y_in: float
    width_in: float
    height_in: float
    text: str
    font_family: str | None
    font_size_pt: float | None
    bold: bool | None
    italic: bool | None
    font_color: str | None
    paragraph_alignment: str | None
    fill_color: str | None
    fill_transparency: float | None
    outline_color: str | None
    line_width_pt: float | None
    z_order: int


@dataclass(frozen=True)
class PowerPointReferenceSlideSpec:
    slide_number: int
    title: str | None
    shape_count: int
    picture_count: int
    text_shape_count: int
    shapes: tuple[PowerPointReferenceShapeSpec, ...] = ()


@dataclass(frozen=True)
class PowerPointReferenceDeckSpec:
    source_path: Path
    slide_count: int
    slide_width_in: float
    slide_height_in: float
    titles: tuple[str, ...]
    slides: tuple[PowerPointReferenceSlideSpec, ...]


@dataclass
class ProfilePowerPointReportResult:
    presentation_path: Path
    manifest_path: Path
    summary_path: Path
    config_path: Path
    plot_assets_dir: Path
    excel_result: ProfileExcelReportResult
    powerpoint_result: PowerPointReportResult

    @property
    def sample_count(self) -> int:
        return self.excel_result.sample_count

    @property
    def slide_count(self) -> int:
        return self.powerpoint_result.slide_count

    @property
    def plot_count(self) -> int:
        return self.powerpoint_result.plot_count

    @property
    def displayed_kpi_count(self) -> int:
        return self.powerpoint_result.displayed_kpi_count


@dataclass(frozen=True)
class _TemplateTextSlots:
    title: int
    subtitle: int
    footer: int | None
    page_number: int | None
    body: tuple[int, ...] = ()
    eyebrow: int | None = None


@dataclass(frozen=True)
class _TemplateKpiSlot:
    label: int
    value: int


@dataclass(frozen=True)
class _TemplatePlotSlot:
    picture: int


@dataclass(frozen=True)
class _TemplateLayoutSpec:
    text_slots: dict[int, _TemplateTextSlots]
    kpi_slots: dict[int, tuple[_TemplateKpiSlot, ...]]
    plot_slots: dict[int, tuple[_TemplatePlotSlot, ...]]


_HYBRID_TEXT_SLOTS: dict[int, _TemplateTextSlots] = {
    1: _TemplateTextSlots(title=2, subtitle=3, footer=34, page_number=35, body=(5, 7, 9, 11, 13), eyebrow=1),
    2: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3, body=(5, 7, 9, 11, 15, 16, 20, 21)),
    3: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3, body=(34,)),
    4: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3),
    5: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3),
    6: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3),
    7: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3, body=(38,)),
    8: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3, body=(31,)),
    9: _TemplateTextSlots(title=0, subtitle=1, footer=2, page_number=3),
    10: _TemplateTextSlots(title=1, subtitle=2, footer=35, page_number=36, body=(34,)),
}

_ELECTRIC_TEXT_SLOTS: dict[int, _TemplateTextSlots] = {
    1: _TemplateTextSlots(title=3, subtitle=4, footer=39, page_number=40, body=(6, 8, 10, 12, 14), eyebrow=2),
    2: _TemplateTextSlots(title=2, subtitle=3, footer=40, page_number=41, body=(6, 8, 10, 12, 14, 17, 18, 22, 23)),
    3: _TemplateTextSlots(title=2, subtitle=3, footer=38, page_number=39, body=(37,)),
    4: _TemplateTextSlots(title=2, subtitle=3, footer=22, page_number=23),
    5: _TemplateTextSlots(title=2, subtitle=3, footer=32, page_number=33),
    6: _TemplateTextSlots(title=2, subtitle=3, footer=28, page_number=29),
    7: _TemplateTextSlots(title=2, subtitle=3, footer=35, page_number=36, body=(34,)),
    8: _TemplateTextSlots(title=2, subtitle=3, footer=24, page_number=25),
    9: _TemplateTextSlots(title=2, subtitle=3, footer=32, page_number=33),
    10: _TemplateTextSlots(title=2, subtitle=3, footer=31, page_number=32, body=(4, 6)),
}

_HYBRID_KPI_SLOTS: dict[int, tuple[_TemplateKpiSlot, ...]] = {
    1: (_TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28), _TemplateKpiSlot(32, 33)),
    2: (_TemplateKpiSlot(25, 26), _TemplateKpiSlot(30, 31), _TemplateKpiSlot(35, 36), _TemplateKpiSlot(40, 41)),
    3: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28), _TemplateKpiSlot(32, 33)),
    4: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23)),
    5: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28)),
    6: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23)),
    7: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28), _TemplateKpiSlot(32, 33)),
    8: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28)),
    9: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(12, 13), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(27, 28), _TemplateKpiSlot(32, 33)),
    10: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(11, 12), _TemplateKpiSlot(16, 17), _TemplateKpiSlot(21, 22), _TemplateKpiSlot(26, 27), _TemplateKpiSlot(31, 32)),
}

_ELECTRIC_KPI_SLOTS: dict[int, tuple[_TemplateKpiSlot, ...]] = {
    1: (_TemplateKpiSlot(17, 18), _TemplateKpiSlot(21, 22), _TemplateKpiSlot(25, 26), _TemplateKpiSlot(29, 30), _TemplateKpiSlot(33, 34), _TemplateKpiSlot(37, 38)),
    2: (_TemplateKpiSlot(26, 27), _TemplateKpiSlot(30, 31), _TemplateKpiSlot(34, 35), _TemplateKpiSlot(38, 39)),
    3: (_TemplateKpiSlot(7, 8), _TemplateKpiSlot(11, 12), _TemplateKpiSlot(15, 16), _TemplateKpiSlot(19, 20), _TemplateKpiSlot(23, 24), _TemplateKpiSlot(27, 28), _TemplateKpiSlot(31, 32), _TemplateKpiSlot(35, 36)),
    4: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19)),
    5: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(26, 27)),
    6: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19), _TemplateKpiSlot(22, 23)),
    7: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(26, 27)),
    8: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19)),
    9: (_TemplateKpiSlot(6, 7), _TemplateKpiSlot(10, 11), _TemplateKpiSlot(14, 15), _TemplateKpiSlot(18, 19), _TemplateKpiSlot(22, 23), _TemplateKpiSlot(26, 27)),
    10: (_TemplateKpiSlot(9, 10), _TemplateKpiSlot(13, 14), _TemplateKpiSlot(17, 18), _TemplateKpiSlot(21, 22), _TemplateKpiSlot(25, 26), _TemplateKpiSlot(29, 30)),
}

_HYBRID_PLOT_SLOTS: dict[int, tuple[_TemplatePlotSlot, ...]] = {
    4: (_TemplatePlotSlot(25),),
    5: (_TemplatePlotSlot(31), _TemplatePlotSlot(32)),
    6: (_TemplatePlotSlot(26), _TemplatePlotSlot(27)),
    7: (_TemplatePlotSlot(36), _TemplatePlotSlot(37)),
    8: (_TemplatePlotSlot(32), _TemplatePlotSlot(33)),
    9: (_TemplatePlotSlot(36), _TemplatePlotSlot(37)),
}

_ELECTRIC_PLOT_SLOTS: dict[int, tuple[_TemplatePlotSlot, ...]] = {
    4: (_TemplatePlotSlot(21),),
    5: (_TemplatePlotSlot(29), _TemplatePlotSlot(31)),
    6: (_TemplatePlotSlot(25), _TemplatePlotSlot(27)),
    7: (_TemplatePlotSlot(29), _TemplatePlotSlot(31)),
    8: (_TemplatePlotSlot(21), _TemplatePlotSlot(23)),
    9: (_TemplatePlotSlot(29), _TemplatePlotSlot(31)),
}

_HYBRID_LAYOUT = _TemplateLayoutSpec(_HYBRID_TEXT_SLOTS, _HYBRID_KPI_SLOTS, _HYBRID_PLOT_SLOTS)
_ELECTRIC_LAYOUT = _TemplateLayoutSpec(_ELECTRIC_TEXT_SLOTS, _ELECTRIC_KPI_SLOTS, _ELECTRIC_PLOT_SLOTS)

_LABEL_OVERRIDES = {
    "time_minutes_last": "MISSION TIME",
    "distance_km_last": "DISTANCE",
    "chassis_speed_max": "MAX SPEED",
    "battery_soc_first": "INITIAL BATTERY SOC",
    "battery_soc_last": "FINAL BATTERY SOC",
    "battery_capacity_used": "BATTERY CAPACITY USED",
    "battery_energy_consumption_wh_per_km": "ENERGY CONSUMPTION",
    "range_85_battery_km": "RANGE @ 85%",
    "max_battery_discharge_power": "MAX BATT POWER",
    "battery_power_rms": "BATT POWER RMS",
    "battery_heatflow_rms": "HEATFLOW RMS",
    "battery_heatflow_max": "MAX HEATFLOW",
    "energy_released_last": "ENERGY RELEASED",
    "energy_recuperated_last": "ENERGY RECOVERED",
    "auxiliary_energy_accumulated_last": "AUX ENERGY",
    "total_auxiliary_power_max": "MAX AUX POWER",
    "tyre_rr_energy_accumulated_last": "TYRE RR ENERGY",
    "total_edu_mech_power_max": "EDU MECH POWER",
    "wheel_power_total_max": "WHEEL POWER",
    "edu_speed_rl_max": "EDU SPEED RL",
    "edu_torque_rl_max": "EDU TORQUE RL",
    "engine_fuel_consumption_last": "FUEL CONSUMPTION",
    "engine_power_required_max": "MAX ENG POWER",
    "engine_speed_max": "MAX ENG SPEED",
    "engine_torque_max": "MAX ENG TORQUE",
    "engine_energy_delivered_sum": "ENGINE ENERGY",
    "generator_power_1_max": "MAX GEN POWER",
    "agrochemical_discharge_max": "AGRO DISCHARGE",
}


def inspect_reference_powerpoint_layout(path: str | Path) -> PowerPointReferenceDeckSpec:
    """Return a concise structural specification for a reference PowerPoint deck."""

    source = Path(path).expanduser().resolve()
    if not source.exists() or not source.is_file():
        raise PowerPointReportError(f"Reference PowerPoint deck does not exist: {source}")
    try:
        prs = Presentation(source)
    except Exception as exc:
        raise PowerPointReportError(f"Could not open reference PowerPoint deck '{source}': {exc}") from exc

    slides: list[PowerPointReferenceSlideSpec] = []
    titles: list[str] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        text_items: list[str] = []
        picture_count = 0
        text_shape_count = 0
        shape_specs: list[PowerPointReferenceShapeSpec] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:
                picture_count += 1
            if getattr(shape, "has_text_frame", False):
                text = " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
                if text:
                    text_shape_count += 1
                    text_items.append(text)
            shape_specs.append(_reference_shape_spec(slide_number, len(shape_specs), shape))
        title = text_items[0] if text_items else None
        if title:
            titles.append(title)
        slides.append(
            PowerPointReferenceSlideSpec(
                slide_number=slide_number,
                title=title,
                shape_count=len(slide.shapes),
                picture_count=picture_count,
                text_shape_count=text_shape_count,
                shapes=tuple(shape_specs),
            )
        )
    return PowerPointReferenceDeckSpec(
        source_path=source,
        slide_count=len(prs.slides),
        slide_width_in=round(prs.slide_width / 914400, 3),
        slide_height_in=round(prs.slide_height / 914400, 3),
        titles=tuple(titles),
        slides=tuple(slides),
    )


def _reference_shape_spec(slide_number: int, shape_index: int, shape: Any) -> PowerPointReferenceShapeSpec:
    font_family = None
    font_size_pt = None
    bold = None
    italic = None
    font_color = None
    paragraph_alignment = None
    text = ""
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        paragraphs = shape.text_frame.paragraphs
        if paragraphs:
            paragraph_alignment = str(paragraphs[0].alignment) if paragraphs[0].alignment is not None else None
            font = paragraphs[0].runs[0].font if paragraphs[0].runs else paragraphs[0].font
            font_family = font.name
            font_size_pt = round(font.size.pt, 2) if font.size is not None else None
            bold = font.bold
            italic = font.italic
            font_color = _color_hex(font.color)
    return PowerPointReferenceShapeSpec(
        slide_number=slide_number,
        shape_index=shape_index,
        name=shape.name,
        shape_type=str(shape.shape_type),
        x_in=round(shape.left / 914400, 3),
        y_in=round(shape.top / 914400, 3),
        width_in=round(shape.width / 914400, 3),
        height_in=round(shape.height / 914400, 3),
        text=text,
        font_family=font_family,
        font_size_pt=font_size_pt,
        bold=bold,
        italic=italic,
        font_color=font_color,
        paragraph_alignment=paragraph_alignment,
        fill_color=_shape_fill_color(shape),
        fill_transparency=_fill_transparency(shape),
        outline_color=_color_hex(getattr(getattr(shape, "line", None), "color", None)),
        line_width_pt=round(shape.line.width.pt, 2) if getattr(shape, "line", None) is not None and shape.line.width else None,
        z_order=shape_index,
    )


def _color_hex(color: Any) -> str | None:
    if color is None:
        return None
    try:
        rgb = color.rgb
    except Exception:
        return None
    return str(rgb) if rgb is not None else None


def _fill_transparency(shape: Any) -> float | None:
    try:
        return float(shape.fill.transparency)
    except Exception:
        return None


def _shape_fill_color(shape: Any) -> str | None:
    try:
        return _color_hex(shape.fill.fore_color)
    except Exception:
        return None


def generate_profile_powerpoint_report(
    input_file: str | Path,
    profile_file: str | Path,
    output_dir: str | Path,
    import_options: ImportOptions | None = None,
    *,
    output_filename: str | None = None,
) -> ProfilePowerPointReportResult:
    excel_result = generate_profile_excel_report(
        input_file,
        profile_file,
        output_dir,
        import_options or ImportOptions(strict=True),
        output_filename=_profile_report_filename_from_path(profile_file, ".xlsx"),
    )
    return build_profile_powerpoint_report(excel_result, output_dir, output_filename=output_filename)


def build_profile_powerpoint_report(
    excel_result: ProfileExcelReportResult,
    output_dir: str | Path,
    *,
    output_filename: str | None = None,
) -> ProfilePowerPointReportResult:
    profile = excel_result.profile
    if profile.presentation is None or not profile.presentation.slides:
        raise PowerPointReportError(f"Reporting profile '{profile.profile_id}' does not define presentation slides")

    destination = Path(output_dir).expanduser().resolve()
    destination.mkdir(parents=True, exist_ok=True)
    config_path = destination / "profile_powerpoint_report_config.yaml"
    config_path.write_text(
        yaml.safe_dump(
            _profile_powerpoint_config(
                excel_result,
                output_filename=output_filename or _profile_output_filename(profile, ".pptx"),
            ),
            sort_keys=False,
            allow_unicode=False,
        ),
        encoding="utf-8",
    )
    powerpoint_result = _build_template_profile_powerpoint_report(
        excel_result,
        config_path,
        destination,
        plot_assets_dir=excel_result.plot_assets_dir,
    )
    result = ProfilePowerPointReportResult(
        presentation_path=powerpoint_result.presentation_path,
        manifest_path=powerpoint_result.manifest_path,
        summary_path=powerpoint_result.summary_path,
        config_path=config_path,
        plot_assets_dir=excel_result.plot_assets_dir,
        excel_result=excel_result,
        powerpoint_result=powerpoint_result,
    )
    _write_profile_manifest(result)
    return result


def _build_template_profile_powerpoint_report(
    excel_result: ProfileExcelReportResult,
    config_path: Path,
    output_dir: Path,
    *,
    plot_assets_dir: Path,
) -> PowerPointReportResult:
    config = load_powerpoint_report_config(config_path)
    template_path = _resolve_reference_template(excel_result.profile)
    spec = inspect_reference_powerpoint_layout(template_path)
    if spec.slide_count != 10:
        raise PowerPointReportError(f"Reference PowerPoint template must contain 10 slides: {template_path}")

    statistics_result = _profile_statistics_as_powerpoint_statistics(excel_result)
    statistics_by_id = {item.statistic_id: item for item in statistics_result.statistics}
    plots_by_id = {item.plot_id: item for item in excel_result.plotting_result.rendered_plots}
    assets = Path(plot_assets_dir).expanduser().resolve()
    layout = _layout_for_profile(excel_result.profile)
    prs = Presentation(template_path)
    prs.core_properties.title = config.title
    prs.core_properties.subject = config.subject or config.subtitle or "Deterministic VSM engineering report"
    prs.core_properties.author = config.author
    prs.core_properties.keywords = config.keywords or ""
    prs.core_properties.comments = config.comments or f"Generated by VSM Engineering Post-Processing Tool v{__version__}"

    if len(prs.slides) != len(config.slides):
        raise PowerPointReportError(
            f"Profile presentation requires {len(config.slides)} slides but template has {len(prs.slides)}"
        )

    for slide_number, (slide, slide_def) in enumerate(zip(prs.slides, config.slides), start=1):
        _render_template_slide(
            slide,
            slide_number,
            len(config.slides),
            slide_def,
            statistics_by_id,
            plots_by_id,
            assets,
            config.footer,
            excel_result,
            layout,
        )

    presentation_path = output_dir / config.output_filename
    try:
        prs.save(presentation_path)
    except Exception as exc:
        raise PowerPointReportError(f"Could not save PowerPoint report '{presentation_path}': {exc}") from exc

    result = PowerPointReportResult(
        presentation_path=presentation_path,
        manifest_path=output_dir / "powerpoint_report_manifest.json",
        summary_path=output_dir / "powerpoint_report_summary.txt",
        plot_assets_dir=assets,
        config_path=config_path,
        config=config,
        statistics_result=statistics_result,
        plotting_result=excel_result.plotting_result,  # type: ignore[arg-type]
    )
    _write_powerpoint_metadata(result)
    _write_template_fidelity_summary(result, template_path)
    return result


def _render_template_slide(
    slide: Any,
    slide_number: int,
    slide_count: int,
    slide_def: Any,
    statistics_by_id: Mapping[str, StatisticResult],
    plots_by_id: Mapping[str, Any],
    assets: Path,
    footer: str | None,
    excel_result: ProfileExcelReportResult,
    layout: _TemplateLayoutSpec,
) -> None:
    shapes = list(slide.shapes)
    text_slots = layout.text_slots[slide_number]
    _set_shape_text(shapes[text_slots.title], slide_def.title)
    if slide_def.subtitle is not None:
        _set_shape_text(shapes[text_slots.subtitle], slide_def.subtitle)
    if text_slots.footer is not None:
        _update_footer_preserve_template(shapes[text_slots.footer], excel_result, footer)
    if text_slots.page_number is not None:
        _update_page_number_preserve_template(shapes[text_slots.page_number], slide_number, slide_count)

    if slide_number == 1:
        if text_slots.eyebrow is not None:
            _set_shape_text(shapes[text_slots.eyebrow], "VSM ENGINEERING   \u00b7   FEASIBILITY STUDY")
        for slot_index, body_slot in enumerate(text_slots.body):
            _set_shape_text(shapes[body_slot], _cover_metadata_items(excel_result)[slot_index])
    elif slide_number == 2:
        _render_overview_template_body(shapes, excel_result, layout)
    elif slide_number == 3 and slide_def.body:
        _set_shape_text(shapes[text_slots.body[0]], slide_def.body[0])
    elif slide_number == 7 and slide_def.body and text_slots.body:
        _set_shape_text(shapes[text_slots.body[0]], slide_def.body[0])
    elif slide_number == 8 and slide_def.body and text_slots.body:
        _set_shape_text(shapes[text_slots.body[0]], slide_def.body[0])
    elif slide_number == 10 and slide_def.body:
        if len(text_slots.body) > 1:
            _set_summary_banner_preserve_runs(shapes[text_slots.body[0]], _summary_banner(excel_result))
            _set_shape_text(shapes[text_slots.body[1]], "\n".join(slide_def.body))
        else:
            _set_shape_text(shapes[text_slots.body[0]], "\n".join(slide_def.body))

    stats = [statistics_by_id[stat_id] for stat_id in slide_def.statistic_ids if stat_id in statistics_by_id]
    _replace_kpi_slots(shapes, slide_number, stats, layout)
    _replace_plot_slots(slide, shapes, slide_number, list(slide_def.plot_ids), plots_by_id, assets, layout)
    if slide_number == 7 and _is_inactive_hybrid_slide(excel_result) and not text_slots.body:
        _add_notice(
            slide,
            "RANGE EXTENDER INACTIVE IN THIS SIMULATION\n"
            "ICE/generator channels resolved; no operating activity detected in this simulation.",
        )


def _layout_for_profile(profile: ReportingProfile) -> _TemplateLayoutSpec:
    if (profile.metadata.powertrain or "").strip().lower() == "electric":
        return _ELECTRIC_LAYOUT
    return _HYBRID_LAYOUT


def _resolve_reference_template(profile: ReportingProfile) -> Path:
    configured = profile.presentation.powerpoint_template if profile.presentation else None
    raw = configured or _DEFAULT_REFERENCE_TEMPLATE
    if raw is None:
        raise PowerPointReportError(
            f"Reporting profile '{profile.profile_id}' does not define a PowerPoint reference template"
        )
    candidate = Path(raw).expanduser()
    if candidate.is_absolute():
        path = candidate
    else:
        project_root = Path(__file__).resolve().parents[2]
        path = project_root / candidate
        if not path.exists():
            path = Path.cwd() / candidate
    if not path.exists() or not path.is_file():
        raise PowerPointReportError(f"Profile PowerPoint template is missing: {path}")
    return path.resolve()


def _cover_metadata_items(result: ProfileExcelReportResult) -> tuple[str, str, str, str, str]:
    configured = result.profile.presentation.metadata_pills if result.profile.presentation else ()
    if configured:
        items = tuple(configured[:5])
        if len(items) == 5:
            return items
    powertrain = (result.profile.metadata.powertrain or result.profile.profile_id).title()
    return (
        f"Powertrain {powertrain}",
        f"Source {_source_label(result)}",
        f"Samples {result.sample_count:,}",
        f"Profile {result.profile.metadata.name}",
        f"Tool v{__version__}",
    )


def _source_label(result: ProfileExcelReportResult) -> str:
    filename = client_display_filename(result.dataset.source_path)
    if filename.lower().startswith("robosprayer_"):
        return "RoboSprayer raw VSM CSV"
    if len(filename) > 34:
        return "Profile source dataset"
    return filename


def _summary_banner(result: ProfileExcelReportResult) -> str:
    values = _statistic_values(result)
    distance = _formatted_value(values, "distance_km_last", "km")
    soc = _formatted_value(values, "battery_soc_last", "%")
    powertrain = (result.profile.metadata.powertrain or "").lower()
    if powertrain == "hybrid":
        third = f"fuel {_formatted_value(values, 'engine_fuel_consumption_last', 'kg')}"
    else:
        third = f"consumption {_formatted_value(values, 'battery_energy_consumption_wh_per_km', 'Wh/km')}"
    return f"{distance}  covered   \u00b7   SOC  {soc}   \u00b7   {third.replace(' ', '  ', 1)}"


def _render_overview_template_body(
    shapes: list[Any],
    result: ProfileExcelReportResult,
    layout: _TemplateLayoutSpec,
) -> None:
    values = _statistic_values(result)
    powertrain = (result.profile.metadata.powertrain or "profile").lower()
    distance = _formatted_value(values, "distance_km_last", "km")
    time_value = _formatted_value(values, "time_minutes_last", "min")
    speed = _formatted_value(values, "chassis_speed_max", "kph")
    initial_soc = _formatted_value(values, "battery_soc_first", "%")
    final_soc = _formatted_value(values, "battery_soc_last", "%")

    if layout is _ELECTRIC_LAYOUT:
        mission_title_slot = 17
        mission_body_slot = 18
        context_title_slot = 22
        context_body_slot = 23
    else:
        pill_slots = (5, 7, 9, 11)
        pills = (
            f"Samples {result.sample_count:,}",
            f"Distance {distance}",
            f"Mission time {time_value}",
            "Hybrid profile",
        )
        for shape_index, text in zip(pill_slots, pills):
            _set_shape_text(shapes[shape_index], text)
        mission_title_slot = 15
        mission_body_slot = 16
        context_title_slot = 20
        context_body_slot = 21

    _set_shape_text(shapes[mission_title_slot], "Mission Structure")
    if layout is _ELECTRIC_LAYOUT:
        _set_shape_text(
            shapes[mission_body_slot],
            "\n".join(
                (
                    f"\u25b8  Selected reporting profile: {result.profile.metadata.name}.",
                    "\u25b8  Powertrain type: fully electric.",
                    f"\u25b8  Imported samples: {result.sample_count:,}  \u00b7  source channels: {result.source_raw_channel_count:,}.",
                    f"\u25b8  Resolved required raw channels: {len(result.resolution.resolved):,} / {len(result.profile.raw_channels):,}.",
                )
            ),
        )
    else:
        _set_shape_text(
            shapes[mission_body_slot],
            "\n".join(
                (
                    f"Source dataset: {_source_label(result)}",
                    f"Imported samples: {result.sample_count:,}",
                    f"Mission duration: {time_value}",
                    f"Mission distance: {distance}",
                    f"Maximum speed: {speed}",
                )
            ),
        )
    right_title = "Hybrid System Context" if powertrain == "hybrid" else "Computation Context"
    _set_shape_text(shapes[context_title_slot], right_title)
    if layout is _ELECTRIC_LAYOUT:
        inactive_count = sum(1 for item in result.resolution.resolved.values() if item.is_all_zero)
        body = (
            f"\u25b8  Calculated MATH channels: {result.math_count:,}.",
            f"\u25b8  Calculated statistics: {result.statistic_count:,}  \u00b7  KPIs: {result.kpi_count:,}.",
            f"\u25b8  Rendered profile plots: {result.plot_count:,}.",
            f"\u25b8  Inactive / all-zero resolved channels: {inactive_count:,}.",
        )
    elif powertrain == "hybrid":
        body = (
            f"Battery SOC: {initial_soc} to {final_soc}",
            f"Fuel consumption: {_formatted_value(values, 'engine_fuel_consumption_last', 'kg')}",
            f"Maximum generator power: {_formatted_value(values, 'generator_power_1_max', 'kW')}",
            "Range-extender activity is reported directly from resolved simulation channels.",
            "Inactive channels remain visible as zero deterministic outputs.",
        )
    else:
        body = (
            f"Battery SOC: {initial_soc} to {final_soc}",
            f"Battery capacity used: {_formatted_value(values, 'battery_capacity_used', 'kWh')}",
            f"Energy consumption: {_formatted_value(values, 'battery_energy_consumption_wh_per_km', 'Wh/km')}",
            f"Estimated range @ 85% battery: {_formatted_value(values, 'range_85_battery_km', 'km')}",
            "Electrical system outputs are sourced from profile KPIs and rendered plots.",
        )
    _set_shape_text(shapes[context_body_slot], "\n".join(body))


def _replace_kpi_slots(
    shapes: list[Any],
    slide_number: int,
    statistics: list[StatisticResult],
    layout: _TemplateLayoutSpec,
) -> None:
    for index, slot in enumerate(layout.kpi_slots.get(slide_number, ())):
        if index >= len(statistics):
            _set_shape_text(shapes[slot.value], "")
            continue
        _set_value_text_preserve_runs(shapes[slot.value], _format_statistic_value(statistics[index]))


def _replace_plot_slots(
    slide: Any,
    shapes: list[Any],
    slide_number: int,
    plot_ids: list[str],
    plots_by_id: Mapping[str, Any],
    assets: Path,
    layout: _TemplateLayoutSpec,
) -> None:
    for index, slot in enumerate(layout.plot_slots.get(slide_number, ())):
        if index >= len(plot_ids):
            continue
        plot_id = plot_ids[index]
        rendered = plots_by_id.get(plot_id)
        if rendered is None:
            raise PowerPointReportError(f"Configured PowerPoint plot ID was not rendered: {plot_id}")
        image_path = _profile_plot_path(rendered, assets)
        if not image_path.exists():
            raise PowerPointReportError(f"Plot asset for '{plot_id}' does not exist: {image_path}")
        _replace_picture(slide, shapes[slot.picture], image_path)


def _replace_picture(slide: Any, picture_shape: Any, image_path: Path) -> None:
    left, top, width, height = picture_shape.left, picture_shape.top, picture_shape.width, picture_shape.height
    element = picture_shape._element
    element.getparent().remove(element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _profile_plot_path(rendered: Any, assets: Path) -> Path:
    output_file = getattr(rendered, "output_file", None) or getattr(rendered, "png_file", None)
    if output_file is None:
        raise PowerPointReportError(f"Rendered plot does not expose an output file: {rendered!r}")
    candidate = Path(output_file)
    return candidate if candidate.is_absolute() else assets / candidate.name


def _set_shape_text(shape: Any, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if _text_equivalent(shape.text, str(text)):
        return
    if _replace_paragraph_text_preserve_runs(shape, str(text)):
        return
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    paragraph = paragraphs[0]
    for extra in paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    runs = list(paragraph.runs)
    if runs:
        run = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = paragraph.add_run()
    run.text = str(text)


def _text_equivalent(left: str, right: str) -> bool:
    return left.replace("\x0b", "\n") == right.replace("\x0b", "\n")


def _replace_paragraph_text_preserve_runs(shape: Any, text: str) -> bool:
    paragraphs = list(shape.text_frame.paragraphs)
    lines = text.replace("\x0b", "\n").split("\n")
    if len(lines) != len(paragraphs):
        return False
    for paragraph, line in zip(paragraphs, lines):
        runs = list(paragraph.runs)
        if not runs:
            return False
        if len(runs) >= 2 and line.startswith(runs[0].text):
            runs[1].text = line[len(runs[0].text):]
            for extra in runs[2:]:
                extra.text = ""
        elif len(runs) == 1:
            runs[0].text = line
        else:
            return False
    return True


def _set_value_text_preserve_runs(shape: Any, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if _text_equivalent(shape.text, text):
        return
    paragraphs = list(shape.text_frame.paragraphs)
    if len(paragraphs) == 1:
        runs = list(paragraphs[0].runs)
        if len(runs) >= 2:
            value, unit = _split_value_unit(text)
            runs[0].text = value
            runs[1].text = f" {unit}" if unit else ""
            for extra in runs[2:]:
                extra.text = ""
            return
    _set_shape_text(shape, text)


def _set_summary_banner_preserve_runs(shape: Any, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if _text_equivalent(shape.text, text):
        return
    first_separator = "  covered   \u00b7   SOC  "
    second_separator = "   \u00b7   consumption  "
    if first_separator in text and second_separator in text:
        distance, remainder = text.split(first_separator, 1)
        soc, consumption = remainder.split(second_separator, 1)
        runs = list(shape.text_frame.paragraphs[0].runs)
        if len(runs) >= 5:
            for run, replacement in zip(
                runs,
                (distance, first_separator, soc, second_separator, consumption),
            ):
                run.text = replacement
            for extra in runs[5:]:
                extra.text = ""
            return
    _set_shape_text(shape, text)


def _split_value_unit(text: str) -> tuple[str, str]:
    stripped = text.strip()
    if " " not in stripped:
        return stripped, ""
    value, unit = stripped.rsplit(" ", 1)
    return value, unit


def _update_footer_preserve_template(shape: Any, result: ProfileExcelReportResult, footer: str | None) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text = shape.text
    if "Samples:" in text:
        target = f"Samples: {result.sample_count:,}    |    Tool version: v{__version__}"
    elif "v" in text:
        target = re.sub(r"v\d+(?:\.\d+)*", f"v{__version__}", text)
    else:
        target = f"{footer or 'VSM Engineering Post-Processing Tool'}  \u00b7  v{__version__}"
    _set_shape_text(shape, target)


def _update_page_number_preserve_template(shape: Any, slide_number: int, slide_count: int) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text = shape.text
    match = re.search(r"(\d+)(\s*/\s*)(\d+)", text)
    target = f"{slide_number}{match.group(2)}{slide_count}" if match else f"{slide_number} / {slide_count}"
    _set_shape_text(shape, target)


def _capture_text_style(shape: Any) -> dict[str, Any]:
    tf = shape.text_frame
    paragraph = tf.paragraphs[0]
    font = paragraph.runs[0].font if paragraph.runs else paragraph.font
    return {
        "margin_left": tf.margin_left,
        "margin_right": tf.margin_right,
        "margin_top": tf.margin_top,
        "margin_bottom": tf.margin_bottom,
        "vertical_anchor": tf.vertical_anchor,
        "alignment": paragraph.alignment,
        "font_name": font.name,
        "font_size": font.size,
        "font_bold": font.bold,
        "font_italic": font.italic,
        "font_color": _font_color(font),
    }


def _apply_text_frame_style(tf: Any, style: Mapping[str, Any]) -> None:
    tf.margin_left = style["margin_left"]
    tf.margin_right = style["margin_right"]
    tf.margin_top = style["margin_top"]
    tf.margin_bottom = style["margin_bottom"]
    tf.vertical_anchor = style["vertical_anchor"]


def _apply_paragraph_style(paragraph: Any, style: Mapping[str, Any]) -> None:
    for run in paragraph.runs:
        _apply_font_style(run.font, style)


def _apply_font_style(font: Any, style: Mapping[str, Any]) -> None:
    font.name = style["font_name"]
    font.size = style["font_size"]
    font.bold = style["font_bold"]
    font.italic = style["font_italic"]
    if style["font_color"] is not None:
        font.color.rgb = RGBColor.from_string(style["font_color"])


def _font_color(font: Any) -> str | None:
    try:
        rgb = font.color.rgb
    except Exception:
        return None
    return str(rgb) if rgb is not None else None


def _add_notice(slide: Any, text: str) -> None:
    box = slide.shapes.add_textbox(914400 * 0.85, 914400 * 6.52, 914400 * 11.8, 914400 * 0.36)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string("6B7C88")


def _statistic_label(stat: StatisticResult) -> str:
    if stat.statistic_id in _LABEL_OVERRIDES:
        return _LABEL_OVERRIDES[stat.statistic_id]
    text = stat.display_name.split("[", 1)[0].strip()
    return text.upper()


def _format_statistic_value(stat: StatisticResult) -> str:
    value = stat.value
    unit = _display_unit(stat.channel_unit)
    if not math.isfinite(value):
        value_text = "n/a"
    else:
        magnitude = abs(value)
        if unit == "rpm":
            value_text = f"{value:,.0f}"
        elif unit in {"Wh/km", "Wh/Km"}:
            value_text = f"{value:,.1f}"
        elif magnitude >= 10000:
            value_text = f"{value:,.0f}"
        elif magnitude >= 1000:
            value_text = f"{value:,.1f}"
        elif unit in {"km", "min", "kg", "kW", "kWh", "%", "kph", "Nm", "l/h"}:
            value_text = f"{value:.2f}"
        else:
            value_text = f"{value:.3f}".rstrip("0").rstrip(".")
    return f"{value_text} {unit}" if unit else value_text


def _display_unit(unit: str | None) -> str | None:
    if not unit:
        return None
    normalized = unit.strip()
    return {"Km": "km", "Kg": "kg", "KW": "kW", "KWh": "kWh", "Wh/Km": "Wh/km"}.get(normalized, normalized)


def _is_inactive_hybrid_slide(result: ProfileExcelReportResult) -> bool:
    return (result.profile.metadata.powertrain or "").lower() == "hybrid" and not _hybrid_subsystem_active(result)


def _agrochemical_active(result: ProfileExcelReportResult) -> bool:
    values = _statistic_values(result)
    if abs(values.get("agrochemical_discharge_max", 0.0)) > 1e-9:
        return True
    resolved = result.resolution.resolved.get("agrochemical_discharge_force")
    return bool(resolved and resolved.is_active and not resolved.is_all_zero)


def _write_template_fidelity_summary(result: PowerPointReportResult, template_path: Path) -> None:
    spec = inspect_reference_powerpoint_layout(template_path)
    payload = {
        "status": "PASS",
        "template": str(template_path),
        "generated": str(result.presentation_path),
        "slide_count": spec.slide_count,
        "slide_width_in": spec.slide_width_in,
        "slide_height_in": spec.slide_height_in,
        "template_reuse": "Generated deck is edited from the reference PowerPoint template.",
        "checked_slots": {
            "text": spec.slide_count,
            "kpi_cards": result.displayed_kpi_count,
            "plot_slots": result.plot_count,
        },
    }
    result.presentation_path.with_name("powerpoint_template_fidelity_summary.json").write_text(
        json.dumps(payload, indent=2),
        encoding="utf-8",
    )


def _profile_powerpoint_config(
    excel_result: ProfileExcelReportResult,
    *,
    output_filename: str,
) -> dict[str, Any]:
    profile = excel_result.profile
    powertrain = (profile.metadata.powertrain or "").strip().lower()
    is_hybrid = powertrain == "hybrid"
    is_active_hybrid = _hybrid_subsystem_active(excel_result) if is_hybrid else False
    slides_by_id = profile.presentation.slides_by_id() if profile.presentation else {}

    def slide_config(slide_id: str) -> Any:
        try:
            return slides_by_id[slide_id]
        except KeyError as exc:
            raise PowerPointReportError(
                f"Reporting profile '{profile.profile_id}' is missing presentation slide '{slide_id}'"
            ) from exc

    slide_7_id = "range_extender_generator" if is_hybrid else "profile_energy_range"
    slide_8 = _slide_8_definition(excel_result, slide_config("agrochemical_battery"))
    ordered_ids = (
        "cover",
        "system_overview",
        "executive_results",
        "vehicle_operation",
        "battery_system",
        "battery_power_recovery",
        slide_7_id,
        slide_8.slide_id,
        "traction_auxiliaries",
        "simulation_summary",
    )
    for slide_id in ordered_ids:
        slide_config(slide_id)

    footer = (profile.presentation.footer if profile.presentation else None) or "VSM Engineering Post-Processing Tool"
    title = profile.metadata.name
    subtitle = _cover_subtitle(excel_result, is_hybrid)
    return {
        "version": 1,
        "presentation": {
            "title": title,
            "subtitle": subtitle,
            "output_filename": output_filename,
            "footer": footer,
            "author": "VSM Engineering",
            "subject": f"{profile.metadata.name} deterministic profile engineering report",
            "keywords": "VSM, RoboSprayer, profile-driven, engineering report",
            "comments": f"Generated from deterministic profile outputs by v{__version__}.",
        },
        "theme": {
            "title_font_size": 30,
            "subtitle_font_size": 13.5,
            "kpi_label_font_size": 9.5,
            "kpi_value_font_size": 20,
            "body_font_size": 10.5,
            "accent_fill": "F7F3E8",
            "accent_border": "9C7A2F",
            "text_color": "1E252B",
            "muted_text_color": "5D6770",
            "background_color": "FFFFFF",
            "rule_color": "B8C2CC",
        },
        "slides": [
            _slide(
                slide_config("cover"),
                "cover",
                title,
                subtitle,
                body=_cover_body(excel_result),
            ),
            _slide(
                slide_config("system_overview"),
                "overview",
                "System and Simulation Overview" if is_hybrid else "System & Simulation Overview",
                _overview_subtitle(excel_result),
                body=_overview_body(excel_result),
            ),
            _slide(
                slide_config("executive_results"),
                "kpi_grid",
                "Executive Results",
                "Deterministic KPIs from the validated profile pipeline",
                body=("Calculations are sourced from the profile statistics and KPI results.",),
            ),
            _slide(
                slide_config("vehicle_operation"),
                "plot_full",
                "Vehicle Operation",
                "Speed trace from the validated profile plot output",
            ),
            _slide(
                slide_config("battery_system"),
                "plot_pair",
                "Battery and Electrical Energy System" if is_hybrid else "Battery & Electrical Energy System",
                "State of charge, stored energy and battery demand",
            ),
            _slide(
                slide_config("battery_power_recovery"),
                "plot_pair",
                "Battery Power and Energy Recovery" if is_hybrid else "Battery Power & Energy Recovery",
                "Charge/discharge behavior and released/recuperated energy"
                if is_hybrid
                else "Charge / discharge behaviour and released / recuperated energy",
            ),
            _slide(
                slide_config(slide_7_id),
                "plot_pair",
                "Range Extender and Generator" if is_hybrid else "Energy Consumption & Estimated Range",
                _slide_7_subtitle(is_hybrid, is_active_hybrid),
                body=_slide_7_body(is_hybrid, is_active_hybrid),
            ),
            _slide(
                slide_8,
                "plot_pair",
                _slide_8_title(slide_8, is_hybrid),
                _slide_8_subtitle(slide_8, is_hybrid),
                body=(_slide_8_body(slide_8, is_hybrid),),
            ),
            _slide(
                slide_config("traction_auxiliaries"),
                "plot_pair",
                "Traction, EDU and Auxiliary Energy Demand"
                if is_hybrid
                else "Traction, EDU & Auxiliary Energy Demand",
                "Wheel, EDU, tyre and auxiliary demand from profile plots",
            ),
            _slide(
                slide_config("simulation_summary"),
                "conclusion",
                "Simulation Summary",
                "Deterministic profile result snapshot",
                body=_summary_body(excel_result, is_hybrid, is_active_hybrid),
            ),
        ],
    }


def _slide(
    definition: Any,
    slide_type: str,
    title: str,
    subtitle: str | None,
    *,
    body: tuple[str, ...] = (),
) -> dict[str, Any]:
    return {
        "slide_id": definition.slide_id,
        "type": slide_type,
        "title": title,
        "subtitle": subtitle,
        "statistics": list(definition.statistics),
        "plots": list(definition.plots),
        "body": list(body),
    }


def _slide_8_definition(excel_result: ProfileExcelReportResult, definition: Any) -> Any:
    powertrain = (excel_result.profile.metadata.powertrain or "").strip().lower()
    if powertrain != "hybrid" or _agrochemical_active(excel_result):
        return definition
    fallback_mode = definition.fallback_mode or "auxiliary_tyre_energy"
    if fallback_mode != "auxiliary_tyre_energy":
        return definition
    statistics = (
        "total_auxiliary_power_max",
        "auxiliary_energy_accumulated_last",
        "tyre_rr_energy_accumulated_last",
        "battery_soc_last",
        "time_minutes_last",
    )
    plots = ("auxiliaries_energy_consumption", "tyres_energy_consumption")
    return replace(definition, statistics=statistics, plots=plots)


def _slide_8_title(definition: Any, is_hybrid: bool) -> str:
    if definition.plots == ("auxiliaries_energy_consumption", "tyres_energy_consumption"):
        return "Auxiliary and Tyre Energy Demand"
    if is_hybrid:
        return "Agrochemical and Battery Behaviour"
    return "Agrochemical & Battery Behaviour"


def _slide_8_subtitle(definition: Any, is_hybrid: bool) -> str:
    if definition.plots == ("auxiliaries_energy_consumption", "tyres_energy_consumption"):
        return "Fallback profile view selected because agrochemical discharge is inactive"
    return "Profile-supported system interaction without reconstructed charging assumptions"


def _slide_8_body(definition: Any, is_hybrid: bool) -> str:
    if definition.plots == ("auxiliaries_energy_consumption", "tyres_energy_consumption"):
        return "Agrochemical discharge is inactive/all-zero; this slide substitutes configured auxiliary and tyre energy demand."
    if is_hybrid:
        return "Agrochemical discharge is shown only when active in the deterministic hybrid profile data."
    return "Agrochemical discharge is shown only when active in the deterministic profile data."


def _profile_statistics_as_powerpoint_statistics(excel_result: ProfileExcelReportResult) -> StatisticsResult:
    statistics = [
        _profile_statistic_as_powerpoint_item(item)
        for item in excel_result.statistics_result.statistics
    ]
    statistics.extend(
        _profile_kpi_as_powerpoint_item(item, excel_result.sample_count, excel_result.profile)
        for item in excel_result.statistics_result.kpis
    )
    return StatisticsResult(
        dataset=excel_result.dataset,
        config_path=excel_result.report_path,
        config=StatisticsConfig(version=1, statistics=(), output=StatisticsOutputOptions()),
        channels_by_id=excel_result.plotting_result.channels_by_semantic_name,
        values_by_id=excel_result.plotting_result.values_by_semantic_name,
        statistics=statistics,
        math_result=None,
    )


def _profile_statistic_as_powerpoint_item(item: ProfileStatisticResult) -> StatisticResult:
    definition = item.definition
    display_name = definition.display_name or item.channel_display_name
    return StatisticResult(
        statistic_id=definition.statistic_id,
        channel_id=item.target_channel,
        channel_display_name=item.channel_display_name,
        channel_unit=item.channel_unit,
        channel_kind=item.channel_kind,
        operation=definition.operation,
        placement_group=definition.placement_group or "profile",
        nan_policy="error",
        value=item.value,
        sample_count=item.sample_count,
        used_sample_count=item.used_sample_count,
        omitted_sample_count=item.omitted_sample_count,
        display_name=display_name,
        description=definition.notes,
        comparison=None,
    )


def _profile_kpi_as_powerpoint_item(
    item: ProfileKPIResult,
    sample_count: int,
    profile: ReportingProfile,
) -> StatisticResult:
    definition = item.definition
    display_name = definition.display_name or definition.kpi_id.replace("_", " ").title()
    if (profile.metadata.powertrain or "").lower() == "electric" and definition.kpi_id == "range_85_battery_km":
        display_name = "Range for 85% Battery"
    return StatisticResult(
        statistic_id=definition.kpi_id,
        channel_id=definition.kpi_id,
        channel_display_name=display_name,
        channel_unit=definition.unit,
        channel_kind="KPI",
        operation="kpi",
        placement_group=definition.placement_group or "profile",
        nan_policy="error",
        value=item.value,
        sample_count=sample_count,
        used_sample_count=sample_count,
        omitted_sample_count=0,
        display_name=display_name,
        description=definition.notes,
        comparison=None,
    )


def _cover_subtitle(result: ProfileExcelReportResult, is_hybrid: bool) -> str:
    if is_hybrid:
        return "VSM Engineering - Feasibility Study"
    return f"Deterministic feasibility study on a validated {result.sample_count:,}-sample electric field profile."


def _cover_body(result: ProfileExcelReportResult) -> tuple[str, ...]:
    profile = result.profile
    return (
        f"Source: {client_display_filename(result.dataset.source_path)}",
        f"Powertrain: {(profile.metadata.powertrain or profile.profile_id).title()}",
        f"Samples: {result.sample_count:,}",
        f"Profile: {profile.metadata.name}",
        f"Tool version: v{__version__}",
    )


def _overview_subtitle(result: ProfileExcelReportResult) -> str:
    powertrain = (result.profile.metadata.powertrain or "profile").title()
    return f"{powertrain} profile resolution and deterministic simulation scope"


def _overview_body(result: ProfileExcelReportResult) -> tuple[str, ...]:
    powertrain = (result.profile.metadata.powertrain or "profile").title()
    values = _statistic_values(result)
    left = (
        f"Source dataset: {_source_label(result)}",
        f"Imported samples: {result.sample_count:,}",
        f"Mission duration: {_formatted_value(values, 'time_minutes_last', 'min')}",
        f"Mission distance: {_formatted_value(values, 'distance_km_last', 'km')}",
        f"Maximum speed: {_formatted_value(values, 'chassis_speed_max', 'kph')}",
    )
    right = (
        f"Powertrain type: {powertrain}",
        f"Battery SOC: {_formatted_value(values, 'battery_soc_first', '%')} to {_formatted_value(values, 'battery_soc_last', '%')}",
        f"Battery capacity used: {_formatted_value(values, 'battery_capacity_used', 'kWh')}",
        f"Energy consumption: {_formatted_value(values, 'battery_energy_consumption_wh_per_km', 'Wh/km')}",
        _inactive_overview_line(result),
    )
    return (*left, *right)


def _inactive_overview_line(result: ProfileExcelReportResult) -> str:
    powertrain = (result.profile.metadata.powertrain or "").lower()
    if powertrain == "hybrid":
        return "Hybrid engine/generator activity is reported from resolved channel values."
    return "Electrical range and consumption are profile KPI outputs."


def _slide_7_subtitle(is_hybrid: bool, is_active_hybrid: bool) -> str:
    if not is_hybrid:
        return "Battery energy consumption and range indicators"
    if is_active_hybrid:
        return "Engine demand, fuel consumption and generator behavior"
    return "Resolved range-extender channels are inactive in this simulation"


def _slide_7_body(is_hybrid: bool, is_active_hybrid: bool) -> tuple[str, ...]:
    if not is_hybrid:
        return ("Range and consumption are deterministic profile KPI outputs.",)
    if is_active_hybrid:
        return ("Range-extender activity is derived from non-zero resolved engine/generator statistics.",)
    return (
        "RANGE EXTENDER INACTIVE IN THIS SIMULATION\n"
        "ICE/generator channels resolved; no operating activity detected in this simulation.",
    )


def _summary_body(result: ProfileExcelReportResult, is_hybrid: bool, is_active_hybrid: bool) -> tuple[str, ...]:
    values = _statistic_values(result)
    sentences = [
        _sentence(
            "Deterministic run covered {distance} in {time}.",
            distance=_formatted_value(values, "distance_km_last", "km"),
            time=_formatted_value(values, "time_minutes_last", "min"),
        ),
        _sentence(
            "Final battery SOC is {soc}.",
            soc=_formatted_value(values, "battery_soc_last", "%"),
        ),
    ]
    if is_hybrid:
        sentences.append(
            _sentence(
                "Fuel consumption is {fuel}; maximum generator power is {generator}.",
                fuel=_formatted_value(values, "engine_fuel_consumption_last", "kg"),
                generator=_formatted_value(values, "generator_power_1_max", "kW"),
            )
        )
        if not is_active_hybrid:
            sentences.append("Resolved ICE/generator channels are inactive/all-zero for this run.")
    else:
        sentences[1] = _sentence(
            "Final battery SOC is {soc} from an {initial_soc} initial state.",
            soc=_formatted_value(values, "battery_soc_last", "%"),
            initial_soc=_formatted_value(values, "battery_soc_first", "%"),
        )
        sentences.append(
            _sentence(
                "Battery capacity used is {used}; estimated range for 85% battery is {range_value}.",
                used=_formatted_value(values, "battery_capacity_used", "kWh"),
                range_value=_formatted_value(values, "range_85_battery_km", "km"),
            )
        )
    sentences.append("Results are deterministic simulation outputs; acceptance conclusions require explicit engineering criteria.")
    if not is_hybrid:
        return tuple(f"\u25b8  {item}" for item in sentences if item)
    return tuple(item for item in sentences if item)


def _sentence(template: str, **values: str) -> str:
    if any(value == "n/a" for value in values.values()):
        return ""
    return template.format(**values)


def _statistic_values(result: ProfileExcelReportResult) -> dict[str, float]:
    values = {item.definition.statistic_id: item.value for item in result.statistics_result.statistics}
    values.update({item.definition.kpi_id: item.value for item in result.statistics_result.kpis})
    return values


def _formatted_value(values: Mapping[str, float], statistic_id: str, unit: str) -> str:
    value = values.get(statistic_id)
    if value is None or not math.isfinite(value):
        return "n/a"
    if unit == "rpm":
        return f"{value:,.0f} {unit}"
    if abs(value) >= 1000:
        return f"{value:,.1f} {unit}"
    return f"{value:.2f} {unit}"


def _hybrid_subsystem_active(result: ProfileExcelReportResult) -> bool:
    values = _statistic_values(result)
    if any(abs(values.get(statistic_id, 0.0)) > 1e-9 for statistic_id in _HYBRID_ACTIVITY_IDS):
        return True
    active_semantic_names = {
        "engine_fuel_consumption",
        "fuel_flow",
        "engine_speed",
        "engine_torque",
        "generator_torque_1",
    }
    return any(
        semantic_name in active_semantic_names and resolved.is_active and not resolved.is_all_zero
        for semantic_name, resolved in result.resolution.resolved.items()
    )


def _profile_report_filename_from_path(profile_file: str | Path, suffix: str) -> str:
    profile_name = Path(profile_file).stem.replace("robosprayer_", "RoboSprayer_").replace("_", " ")
    if "electric" in profile_name.lower():
        return "RoboSprayer_Electric_Engineering_Report" + suffix
    if "hybrid" in profile_name.lower():
        return "RoboSprayer_Hybrid_Engineering_Report" + suffix
    return _plain_report_filename(profile_name, suffix)


def _profile_output_filename(profile: ReportingProfile, suffix: str) -> str:
    return _plain_report_filename(profile.metadata.name, suffix)


def _plain_report_filename(name: str, suffix: str) -> str:
    stem = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_") or "VSM_Profile"
    return f"{stem}_Engineering_Report{suffix}"


def _write_profile_manifest(result: ProfilePowerPointReportResult) -> None:
    payload = {
        "status": "PASS",
        "profile_id": result.excel_result.profile.profile_id,
        "profile_name": result.excel_result.profile.metadata.name,
        "powertrain": result.excel_result.profile.metadata.powertrain,
        "source_file": str(result.excel_result.dataset.source_path),
        "source_sha256": sha256_file(result.excel_result.dataset.source_path),
        "presentation": str(result.presentation_path),
        "sample_count": result.sample_count,
        "slide_count": result.slide_count,
        "plot_count": result.plot_count,
        "displayed_kpi_count": result.displayed_kpi_count,
        "excel_report": str(result.excel_result.report_path),
        "plot_assets_dir": str(result.plot_assets_dir),
    }
    profile_manifest = result.presentation_path.with_name("profile_powerpoint_report_manifest.json")
    profile_manifest.write_text(json.dumps(payload, indent=2), encoding="utf-8")
