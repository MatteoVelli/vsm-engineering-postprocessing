from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from vsm_postprocessing.errors import ConfigurationError, PowerPointReportError
from vsm_postprocessing.importer import ImportOptions
from vsm_postprocessing.powerpoint_report_engine import (
    generate_powerpoint_report,
    load_powerpoint_report_config,
)
from conftest import CAIMAN_REFERENCE_DESCRIPTION, CAIMAN_REFERENCE_XLSX, require_private_reference_file

PROJECT_ROOT = Path(__file__).resolve().parents[1]
POWERPOINT_CONFIG = PROJECT_ROOT / "config" / "powerpoint_report_example.yaml"
STATISTICS_CONFIG = PROJECT_ROOT / "config" / "statistics_excel_report.yaml"
PLOTTING_CONFIG = PROJECT_ROOT / "config" / "plotting_example.yaml"
MATH_CONFIG = PROJECT_ROOT / "config" / "math_channels_example.yaml"


def test_powerpoint_config_loads() -> None:
    config = load_powerpoint_report_config(POWERPOINT_CONFIG)
    assert config.version == 1
    assert config.output_filename == "vsm_engineering_report.pptx"
    assert len(config.slides) == 4
    assert config.slides[1].slide_type == "plot_pair"
    assert config.slides[1].plot_ids == ("speed_vs_time", "battery_soc_vs_distance")


def test_powerpoint_config_rejects_invalid_slide_type(tmp_path: Path) -> None:
    path = tmp_path / "bad.yaml"
    path.write_text(
        """version: 1
presentation:
  title: Test
slides:
  - slide_id: bad
    type: magic
    title: Bad
    statistics: []
    plots: []
""",
        encoding="utf-8",
    )
    with pytest.raises(ConfigurationError, match="type must be one of"):
        load_powerpoint_report_config(path)


def test_powerpoint_config_rejects_three_plots(tmp_path: Path) -> None:
    path = tmp_path / "bad.yaml"
    path.write_text(
        """version: 1
presentation:
  title: Test
slides:
  - slide_id: plots
    type: plot_pair
    title: Too many
    statistics: []
    plots: [a, b, c]
""",
        encoding="utf-8",
    )
    with pytest.raises(ConfigurationError, match="one or two"):
        load_powerpoint_report_config(path)


def test_supplied_source_workbook_powerpoint_acceptance(tmp_path: Path) -> None:
    source_workbook = require_private_reference_file(CAIMAN_REFERENCE_XLSX, CAIMAN_REFERENCE_DESCRIPTION)
    result = generate_powerpoint_report(
        source_workbook,
        POWERPOINT_CONFIG,
        STATISTICS_CONFIG,
        PLOTTING_CONFIG,
        tmp_path / "powerpoint",
        ImportOptions(strict=True),
        math_config_file=MATH_CONFIG,
    )
    assert result.presentation_path.exists()
    assert result.manifest_path.exists()
    assert result.summary_path.exists()
    assert result.slide_count == 4
    assert result.plot_count == 6
    assert result.statistic_count >= 10

    prs = Presentation(result.presentation_path)
    assert len(prs.slides) == 4
    titles = []
    for slide in prs.slides:
        texts = [getattr(shape, "text", "") for shape in slide.shapes]
        titles.append(next(text for text in texts if text.strip()))
    assert titles[0] == "Hybrid SP Caiman – Engineering Summary"
    assert titles[1] == "Caiman SP Hybrid – Source 74 ha Field Cycle"
