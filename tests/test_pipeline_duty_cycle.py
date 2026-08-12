from __future__ import annotations

import json
import zipfile
from pathlib import Path
from xml.etree import ElementTree

import numpy as np
import pytest
from openpyxl import load_workbook
from pptx import Presentation

from vsm_postprocessing.duty_cycle import (
    WorkbookRowProfileProvider,
    compose_duty_cycle,
    export_pipeline_dataset,
    load_duty_cycle_config,
    load_profile_provider_config,
)
from vsm_postprocessing.importer import ImportOptions, load_data_file
from vsm_postprocessing.pipeline_engine import load_pipeline_config, run_pipeline

PROJECT_ROOT = Path(__file__).resolve().parents[1]
SOURCE_WORKBOOK = PROJECT_ROOT / "reference_files" / "Sprayer_Caiman_SP_9300Kg_Hybrid_Gen80kW_30kph_74Ht_4000KgAQ_57-4pcSOC_5-80_1C2G_02.xlsx"
REFERENCE_WORKBOOK = PROJECT_ROOT / "reference_files" / "Sprayer_Caiman_SP_9300Kg_Electrification_03.xlsx"
SCENARIO_CONFIG = PROJECT_ROOT / "config" / "duty_cycle_sergio_reference.yaml"
PROVIDER_CONFIG = PROJECT_ROOT / "config" / "duty_cycle_profiles_sergio_reference.yaml"
PIPELINE_CONFIG = PROJECT_ROOT / "config" / "end_to_end_sergio_duty_cycle.yaml"


@pytest.mark.skipif(
    not SOURCE_WORKBOOK.exists() or not REFERENCE_WORKBOOK.exists(),
    reason="Sergio reference workbooks are not present",
)
def test_pipeline_dataset_export_roundtrips_full_composition(tmp_path: Path) -> None:
    source = load_data_file(SOURCE_WORKBOOK)
    scenario = load_duty_cycle_config(SCENARIO_CONFIG)
    provider = WorkbookRowProfileProvider(load_profile_provider_config(PROVIDER_CONFIG), REFERENCE_WORKBOOK)
    composition = compose_duty_cycle(scenario, source, provider)

    export_path = export_pipeline_dataset(composition, tmp_path / "duty_cycle_dataset.csv")
    reloaded = load_data_file(
        export_path,
        ImportOptions(
            header_row=1,
            unit_row=2,
            data_start_row=3,
            last_channel_column=source.quality.channel_count,
            time_channel=source.channels[source.channel_index(source.quality.time_channel_id)].source_name,
            strict=True,
        ),
    )

    assert reloaded.quality.sample_count == 17418
    assert reloaded.quality.channel_count == 70
    assert [channel.channel_id for channel in reloaded.channels] == [channel.channel_id for channel in source.channels]
    np.testing.assert_allclose(reloaded.values, composition.values, rtol=0.0, atol=0.0)


def test_full_duty_cycle_pipeline_config_loads_optional_composer_block(tmp_path: Path) -> None:
    source = tmp_path / "input.csv"
    source.write_text("Time,Signal\ns,kW\n0,1\n1,2\n", encoding="utf-8")
    scenario = tmp_path / "scenario.yaml"
    scenario.write_text("scenario: {}\n", encoding="utf-8")
    provider = tmp_path / "provider.yaml"
    provider.write_text("provider: {}\n", encoding="utf-8")
    profile = tmp_path / "profile.xlsx"
    profile.write_bytes(b"placeholder")
    config_names = (
        "channel_selection",
        "math_channels",
        "statistics",
        "plotting",
        "excel_statistics",
        "excel_report",
    )
    for name in config_names:
        (tmp_path / f"{name}.yaml").write_text("version: 1\n", encoding="utf-8")
    pipeline = tmp_path / "pipeline.yaml"
    pipeline.write_text(
        """version: 1
input:
  file: input.csv
duty_cycle:
  scenario: scenario.yaml
  profile_provider: provider.yaml
  profile_workbook: profile.xlsx
configs:
  channel_selection: channel_selection.yaml
  math_channels: math_channels.yaml
  statistics: statistics.yaml
  plotting: plotting.yaml
  excel_statistics: excel_statistics.yaml
  excel_report: excel_report.yaml
output:
  root_dir: output
""",
        encoding="utf-8",
    )
    config = load_pipeline_config(pipeline)
    assert config.duty_cycle is not None
    assert config.duty_cycle.scenario_config == scenario.resolve()
    assert config.duty_cycle.profile_provider_config == provider.resolve()
    assert config.duty_cycle.profile_workbook == profile.resolve()


@pytest.mark.skipif(
    not SOURCE_WORKBOOK.exists() or not REFERENCE_WORKBOOK.exists(),
    reason="Sergio reference workbooks are not present",
)
def test_full_duty_cycle_runs_through_normal_reporting_pipeline(tmp_path: Path) -> None:
    raw = PIPELINE_CONFIG.read_text(encoding="utf-8")
    raw = raw.replace(
        "../outputs/end_to_end_sergio_duty_cycle",
        str((tmp_path / "full_mission").resolve()).replace("\\", "/"),
    )
    raw = raw.replace(
        "../reference_files/Sprayer_Caiman_SP_9300Kg_Hybrid_Gen80kW_30kph_74Ht_4000KgAQ_57-4pcSOC_5-80_1C2G_02.xlsx",
        str(SOURCE_WORKBOOK.resolve()).replace("\\", "/"),
    )
    raw = raw.replace(
        "../reference_files/Sprayer_Caiman_SP_9300Kg_Electrification_03.xlsx",
        str(REFERENCE_WORKBOOK.resolve()).replace("\\", "/"),
    )
    for filename in (
        "duty_cycle_sergio_reference.yaml",
        "duty_cycle_profiles_sergio_reference.yaml",
        "channel_selection_example.yaml",
        "math_channels_example.yaml",
        "statistics_excel_report.yaml",
        "plotting_example.yaml",
        "excel_report_duty_cycle.yaml",
        "powerpoint_report_duty_cycle.yaml",
    ):
        raw = raw.replace(filename, str((PROJECT_ROOT / "config" / filename).resolve()).replace("\\", "/"))
    config_path = tmp_path / "full_duty_cycle_pipeline.yaml"
    config_path.write_text(raw, encoding="utf-8")

    result = run_pipeline(config_path)

    assert result.status == "PASS"
    assert result.completed_stage_count == 8
    assert [stage.name for stage in result.stages] == [
        "inspection",
        "duty_cycle",
        "channel_selection",
        "math_channels",
        "statistics",
        "plotting",
        "excel_report",
        "powerpoint_report",
    ]
    assert result.stages[0].metrics["samples"] == 1866
    assert result.stages[1].metrics["samples"] == 17418
    assert result.stages[1].metrics["phases"] == 12
    assert result.stages[1].metrics["external_profile_phases"] == 4
    assert result.stages[1].metrics["final_time_min"] == pytest.approx(290.28333333333336)
    assert result.stages[1].metrics["final_distance_km"] == pytest.approx(114.0011, abs=1e-4)
    assert result.stages[1].metrics["final_soc_pct"] == pytest.approx(23.9383, abs=1e-4)
    assert result.stages[1].metrics["max_speed_kph"] == pytest.approx(62.6233, abs=1e-4)
    assert result.stages[1].metrics["max_generator_kw"] == pytest.approx(80.02669, abs=1e-4)

    for stage in result.stages[2:]:
        assert stage.metrics["samples"] == 17418
    assert result.stages[3].metrics["source_channels"] == 70
    assert result.stages[3].metrics["math_channels"] == 13
    assert result.stages[4].metrics["statistics"] == 53
    assert result.stages[5].metrics["plots"] == 24
    assert result.stages[6].metrics["report_channels"] == 70
    assert result.report_path is not None and result.report_path.exists()
    assert result.powerpoint_path is not None and result.powerpoint_path.exists()
    assert result.processing_input_file.name == "duty_cycle_dataset.csv"
    assert result.processing_input_file.exists()

    manifest = json.loads(result.manifest_path.read_text(encoding="utf-8"))
    assert manifest["duty_cycle"] is not None
    assert manifest["processing_input_file"] == str(result.processing_input_file)
    assert manifest["stage_count"] == 8

    workbook = load_workbook(result.report_path, data_only=True)
    try:
        report = workbook["Report"]
        metadata = workbook["Metadata"]
        assert report.max_column >= 98
        assert report.max_row == 17424
        assert report["A3"].value == "Track_Time"
        assert report["BR3"].value == "Total Generator Power"
        assert report["B17422"].value == pytest.approx(290.28333333333336)
        assert report["G17422"].value == pytest.approx(114.0011, abs=1e-4)
        assert report["R17422"].value == pytest.approx(23.9383, abs=1e-4)
        assert report["AD17422"].value == pytest.approx(39.84212, abs=1e-5)
        assert report["BT3"].value == "Time [min]"
        assert len(report._charts) == 18
        assert len(report._images) == 0
        metadata_values = [
            str(cell.value)
            for row in metadata.iter_rows()
            for cell in row
            if cell.value is not None
        ]
        assert not any("C:\\Users" in value or "Desktop\\Agro" in value for value in metadata_values)
    finally:
        workbook.close()

    assert result.powerpoint_path is not None
    prs = Presentation(result.powerpoint_path)
    assert len(prs.slides) == 10
    assert prs.core_properties.title == "Hybrid SP Caiman Sprayer - Feasibility Study"
    assert prs.core_properties.subject == "Hybrid SP Caiman deterministic duty-cycle feasibility study"
    assert prs.core_properties.author == "VSM Engineering"
    slide_texts = [
        [
            " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and " ".join(paragraph.text for paragraph in shape.text_frame.paragraphs).strip()
        ]
        for slide in prs.slides
    ]
    titles = [texts[0] for texts in slide_texts]
    assert titles == [
        "Hybrid SP Caiman Sprayer - Feasibility Study",
        "System and Mission Overview",
        "Duty-Cycle Executive Results",
        "Duty-Cycle Vehicle Operation",
        "Battery and Electrical Energy System",
        "Battery Power and Energy Recovery",
        "Range Extender and Generator",
        "Opportunity Charging During Loading",
        "Traction, EDU and Auxiliary Energy Demand",
        "Simulation Summary",
    ]
    assert len(titles) == len(set(titles))
    all_text = "\n".join(text for slide in slide_texts for text in slide)
    for expected in ("114.00 km", "290.28 min", "23.94 %", "39.84 kg", "80.03 kW", "1/10", "10/10"):
        assert expected in all_text
    assert "VSM Engineering Post-Processing Tool | v" in all_text
    assert not any("C:\\Users" in text or "Desktop\\Agro" in text for text in all_text.splitlines())
    picture_count = sum(
        1
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "shape_type", None) == 13
    )
    assert picture_count == 11
    assert result.stages[7].metrics["slides"] == 10
    assert result.stages[7].metrics["plots"] == 10

    ppt_manifest = json.loads(
        (result.powerpoint_path.parent / "powerpoint_report_manifest.json").read_text(encoding="utf-8")
    )
    assert ppt_manifest["slide_count"] == 10
    assert ppt_manifest["displayed_kpi_count"] == 58
    assert ppt_manifest["appendix_slide_count"] == 0
    assert ppt_manifest["matplotlib_image_count"] == 10
    assert ppt_manifest["matplotlib_image_placement_count"] == 11
    assert set(ppt_manifest["plot_files"]) == set(ppt_manifest["plots_used"])
    assert all(Path(path).exists() for path in ppt_manifest["plot_files"].values())

    with zipfile.ZipFile(result.powerpoint_path) as archive:
        names = set(archive.namelist())
        assert "[Content_Types].xml" in names
        media = [name for name in names if name.startswith("ppt/media/")]
        assert len(media) == 10
        missing_media: list[tuple[str, str]] = []
        for rel_name in [name for name in names if name.startswith("ppt/slides/_rels/")]:
            root = ElementTree.fromstring(archive.read(rel_name))
            for node in root:
                target = node.attrib.get("Target", "")
                if target.startswith("../media/"):
                    media_name = "ppt/media/" + target.split("../media/", 1)[1]
                    if media_name not in names:
                        missing_media.append((rel_name, target))
        assert missing_media == []
