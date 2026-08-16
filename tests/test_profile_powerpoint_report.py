from __future__ import annotations

import copy
import hashlib
import json
import zipfile
from dataclasses import replace
from pathlib import Path

import pytest
from pptx import Presentation

from vsm_postprocessing.importer import ImportOptions
from vsm_postprocessing.profile_powerpoint_report_engine import (
    build_profile_powerpoint_report,
    _hybrid_subsystem_active,
    _profile_powerpoint_config,
    inspect_reference_powerpoint_layout,
)
from vsm_postprocessing.ui_config import generate_reporting_profile_engineering_report

from conftest import (
    ROBOSPRAYER_REFERENCE_CSV,
    ROBOSPRAYER_REFERENCE_DESCRIPTION,
    require_private_reference_file,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
ELECTRIC_REFERENCE_DECK = PROJECT_ROOT / "reference_files" / "RoboSprayer_Electric_Report_FINAL.pptx"
HYBRID_REFERENCE_DECK = PROJECT_ROOT / "reference_files" / "RoboSprayer_Hybrid_Engineering_Report.pptx"
ELECTRIC_PROFILE = PROJECT_ROOT / "config" / "report_profiles" / "robosprayer_electric.yaml"
HYBRID_PROFILE = PROJECT_ROOT / "config" / "report_profiles" / "robosprayer_hybrid.yaml"


def _robosprayer_csv() -> Path:
    return require_private_reference_file(ROBOSPRAYER_REFERENCE_CSV, ROBOSPRAYER_REFERENCE_DESCRIPTION)


@pytest.fixture(scope="module")
def electric_report(tmp_path_factory: pytest.TempPathFactory):
    return generate_reporting_profile_engineering_report(
        _robosprayer_csv(),
        ELECTRIC_PROFILE,
        tmp_path_factory.mktemp("profile_ppt_electric"),
        ImportOptions(strict=True),
    )


@pytest.fixture(scope="module")
def hybrid_report(tmp_path_factory: pytest.TempPathFactory):
    return generate_reporting_profile_engineering_report(
        _robosprayer_csv(),
        HYBRID_PROFILE,
        tmp_path_factory.mktemp("profile_ppt_hybrid"),
        ImportOptions(strict=True),
    )


def test_reference_powerpoint_layout_inspection_matches_final_profile_decks() -> None:
    assert ELECTRIC_REFERENCE_DECK.exists()
    assert HYBRID_REFERENCE_DECK.exists()
    electric = inspect_reference_powerpoint_layout(ELECTRIC_REFERENCE_DECK)
    hybrid = inspect_reference_powerpoint_layout(HYBRID_REFERENCE_DECK)

    assert electric.slide_count == hybrid.slide_count == 10
    assert electric.slide_width_in == hybrid.slide_width_in == pytest.approx(13.333)
    assert electric.slide_height_in == hybrid.slide_height_in == pytest.approx(7.5)
    assert _reference_titles(electric.source_path) == [
        "RoboSprayer Electric",
        "System & Simulation Overview",
        "Executive Results",
        "Vehicle Operation",
        "Battery & Electrical Energy System",
        "Battery Power & Energy Recovery",
        "Energy Consumption & Estimated Range",
        "Agrochemical & Battery Behaviour",
        "Traction, EDU & Auxiliary Energy Demand",
        "Simulation Summary",
    ]
    assert _reference_titles(hybrid.source_path) == [
        "RoboSprayer Hybrid",
        "System and Simulation Overview",
        "Executive Results",
        "Vehicle Operation",
        "Battery and Electrical Energy System",
        "Battery Power and Energy Recovery",
        "Range Extender and Generator",
        "Auxiliary and Tyre Energy Demand",
        "Traction, EDU and Auxiliary Energy Demand",
        "Simulation Summary",
    ]
    assert electric.slides[3].shapes[21].x_in == pytest.approx(3.157)
    assert hybrid.slides[3].shapes[25].x_in == pytest.approx(3.419)


def test_profile_powerpoint_generates_reopenable_electric_and_hybrid_decks(electric_report, hybrid_report) -> None:
    assert electric_report.presentation_path.name == "RoboSprayer_Electric_Engineering_Report.pptx"
    assert hybrid_report.presentation_path.name == "RoboSprayer_Hybrid_Engineering_Report.pptx"

    for result in (electric_report, hybrid_report):
        assert result.report_path.exists()
        assert result.presentation_path.exists()
        assert result.slide_count == 10
        with zipfile.ZipFile(result.presentation_path) as package:
            assert package.testzip() is None
            media = [name for name in package.namelist() if name.startswith("ppt/media/")]
        assert len(media) >= 10
        prs = Presentation(result.presentation_path)
        assert len(prs.slides) == 10


def test_profile_powerpoint_slide_titles_are_profile_conditional(electric_report, hybrid_report) -> None:
    assert _slide_titles(electric_report.presentation_path) == [
        "RoboSprayer Electric",
        "System & Simulation Overview",
        "Executive Results",
        "Vehicle Operation",
        "Battery & Electrical Energy System",
        "Battery Power & Energy Recovery",
        "Energy Consumption & Estimated Range",
        "Agrochemical & Battery Behaviour",
        "Traction, EDU & Auxiliary Energy Demand",
        "Simulation Summary",
    ]
    assert _slide_titles(hybrid_report.presentation_path) == [
        "RoboSprayer Hybrid",
        "System and Simulation Overview",
        "Executive Results",
        "Vehicle Operation",
        "Battery and Electrical Energy System",
        "Battery Power and Energy Recovery",
        "Range Extender and Generator",
        "Auxiliary and Tyre Energy Demand",
        "Traction, EDU and Auxiliary Energy Demand",
        "Simulation Summary",
    ]


def test_profile_powerpoint_uses_profile_kpis_and_omits_caiman_leakage(electric_report) -> None:
    text = _visible_text(electric_report.presentation_path)
    values = {
        item.definition.statistic_id: item.value
        for item in electric_report.excel_result.statistics_result.statistics
    }
    values.update(
        {
            item.definition.kpi_id: item.value
            for item in electric_report.excel_result.statistics_result.kpis
        }
    )

    assert f"{values['battery_soc_last']:.2f} %" in text
    assert f"{values['battery_power_rms']:.2f} kW" in text
    assert f"{values['distance_km_last']:.2f} km" in text
    assert f"{values['chassis_speed_max']:.2f} kph" in text
    assert "Agrochemical & Battery Behaviour" in text
    assert "Auxiliary and Tyre Energy Demand" not in text
    assert "Generator" not in text
    assert "Fuel" not in text
    assert "Caiman" not in text
    assert "17,418" not in text
    assert "114.00" not in text
    assert "290.28" not in text
    assert "39.84" not in text
    assert "23.94" not in text
    assert "C:\\Users\\" not in text
    assert "Desktop\\Agro Project" not in text


def test_profile_powerpoint_hybrid_inactive_subsystem_handling(hybrid_report) -> None:
    text = _visible_text(hybrid_report.presentation_path)

    assert _hybrid_subsystem_active(hybrid_report.excel_result) is False
    assert "Range Extender and Generator" in text
    assert "RANGE EXTENDER INACTIVE IN THIS SIMULATION" in text
    assert "ICE/generator channels resolved; no operating activity detected in this simulation." in text
    assert "Resolved ICE/generator channels are inactive/all-zero for this run." in text


def test_profile_powerpoint_future_active_hybrid_uses_active_slide_copy(hybrid_report) -> None:
    excel_result = copy.copy(hybrid_report.excel_result)
    statistics_result = copy.copy(hybrid_report.excel_result.statistics_result)
    statistics_result.statistics = [
        replace(item, value=1200.0)
        if item.definition.statistic_id == "engine_speed_max"
        else item
        for item in statistics_result.statistics
    ]
    excel_result.statistics_result = statistics_result

    assert _hybrid_subsystem_active(excel_result) is True
    config = _profile_powerpoint_config(excel_result, output_filename="active_hybrid.pptx")
    slide_7 = config["slides"][6]
    assert slide_7["title"] == "Range Extender and Generator"
    assert slide_7["subtitle"] == "Engine demand, fuel consumption and generator behavior"
    assert slide_7["body"] == [
        "Range-extender activity is derived from non-zero resolved engine/generator statistics."
    ]


def test_profile_powerpoint_plot_ids_and_manifest_are_deterministic(electric_report, hybrid_report) -> None:
    electric_manifest = json.loads(electric_report.manifest_path.read_text(encoding="utf-8"))
    hybrid_manifest = json.loads(hybrid_report.manifest_path.read_text(encoding="utf-8"))

    assert electric_manifest["slides"][3]["plots"] == ["speed_vs_distance"]
    assert electric_manifest["slides"][6]["plots"] == [
        "battery_energy_distance_based",
        "auxiliaries_energy_consumption",
    ]
    assert hybrid_manifest["slides"][6]["plots"] == [
        "generator_power",
        "engine_power_and_fuel_vs_time",
    ]
    assert electric_manifest["slides"][7]["plots"] == [
        "agrochemical_discharge_vs_distance",
        "agrochemical_discharge_and_charge_vs_time",
    ]
    assert electric_manifest["slides"][6]["title"] == "Energy Consumption & Estimated Range"
    assert hybrid_manifest["slides"][6]["title"] == "Range Extender and Generator"
    assert electric_manifest["slides"][7]["title"] == "Agrochemical & Battery Behaviour"


def test_profile_powerpoint_preserves_electric_reference_geometry_and_style(electric_report) -> None:
    generated = Presentation(electric_report.presentation_path)
    reference = Presentation(ELECTRIC_REFERENCE_DECK)

    assert generated.slide_width == reference.slide_width
    assert generated.slide_height == reference.slide_height
    assert len(generated.slides) == len(reference.slides) == 10

    generated_cover = list(generated.slides[0].shapes)
    reference_cover = list(reference.slides[0].shapes)
    assert generated_cover[3].left == reference_cover[3].left
    assert generated_cover[3].top == reference_cover[3].top
    assert generated_cover[3].width == reference_cover[3].width
    assert _first_run_font(generated_cover[3]).name == "Cambria"
    assert str(_first_run_font(generated_cover[3]).color.rgb) == "FFFFFF"

    generated_slide4 = list(generated.slides[3].shapes)
    reference_slide4 = list(reference.slides[3].shapes)
    assert _shape_geometry(reference_slide4[21]) in _picture_geometries(generated_slide4)


def test_profile_powerpoint_preserves_final_reference_text_runs(electric_report, hybrid_report) -> None:
    electric = Presentation(electric_report.presentation_path)
    electric_reference = Presentation(ELECTRIC_REFERENCE_DECK)
    hybrid = Presentation(hybrid_report.presentation_path)
    hybrid_reference = Presentation(HYBRID_REFERENCE_DECK)

    electric_cover = list(electric.slides[0].shapes)
    assert [electric_cover[index].text for index in (6, 8, 10, 12, 14)] == [
        "GVW 3,500 kg",
        "Electric powertrain",
        "12 kph field",
        "Battery 50 kW",
        "Motor 63 rpm",
    ]
    assert "Source RoboSprayer raw VSM CSV" not in _visible_text(electric_report.presentation_path)
    assert _run_texts(electric_cover[18]) == ["12.00", " km"]
    assert _run_styles(electric_cover[18]) == _run_styles(list(electric_reference.slides[0].shapes)[18])

    slide_2_body = list(electric.slides[1].shapes)[18]
    assert len(slide_2_body.text_frame.paragraphs) == 4
    for paragraph in slide_2_body.text_frame.paragraphs:
        assert _run_texts_from_paragraph(paragraph)[0] == "\u25b8  "

    slide_10_banner = list(electric.slides[9].shapes)[4]
    assert _run_texts(slide_10_banner) == [
        "12.00 km",
        "  covered   \u00b7   SOC  ",
        "12.69 %",
        "   \u00b7   consumption  ",
        "2,804.6 Wh/km",
    ]
    assert len({style[3] for style in _run_styles(slide_10_banner)}) >= 3
    assert list(electric.slides[9].shapes)[31].text == list(electric_reference.slides[9].shapes)[31].text
    assert list(hybrid.slides[9].shapes)[35].text == list(hybrid_reference.slides[9].shapes)[35].text

    hybrid_notice = _find_text_shape(hybrid.slides[6], "RANGE EXTENDER INACTIVE")
    reference_notice = _find_text_shape(hybrid_reference.slides[6], "RANGE EXTENDER INACTIVE")
    assert hybrid_notice.text == reference_notice.text
    assert len(hybrid_notice.text_frame.paragraphs) == 1
    assert _run_texts(hybrid_notice) == _run_texts(reference_notice)


def test_profile_powerpoint_preserves_hybrid_reference_plot_slots(hybrid_report) -> None:
    generated = Presentation(hybrid_report.presentation_path)
    reference = Presentation(HYBRID_REFERENCE_DECK)

    assert generated.slide_width == reference.slide_width
    assert generated.slide_height == reference.slide_height
    assert len(generated.slides) == len(reference.slides) == 10

    generated_slide7 = list(generated.slides[6].shapes)
    reference_slide7 = list(reference.slides[6].shapes)
    for index in (36, 37):
        assert _shape_geometry(reference_slide7[index]) in _picture_geometries(generated_slide7)


def test_profile_powerpoint_embedded_media_matches_final_references(electric_report, hybrid_report) -> None:
    assert _media_hashes(electric_report.presentation_path) == _media_hashes(ELECTRIC_REFERENCE_DECK)
    assert _media_hashes(hybrid_report.presentation_path) == _media_hashes(HYBRID_REFERENCE_DECK)


def test_profile_powerpoint_dynamic_values_update_without_losing_reference_runs(
    electric_report,
    hybrid_report,
    tmp_path: Path,
) -> None:
    electric_result = _with_values(
        electric_report.excel_result,
        {
            "distance_km_last": 13.5,
            "battery_soc_last": 44.4,
            "battery_energy_consumption_wh_per_km": 1234.5,
        },
    )
    electric_output = build_profile_powerpoint_report(electric_result, tmp_path / "electric_dynamic")
    electric = Presentation(electric_output.presentation_path)
    assert _run_texts(list(electric.slides[0].shapes)[18]) == ["13.50", " km"]
    assert _run_texts(list(electric.slides[0].shapes)[26]) == ["44.40", " %"]
    assert _run_texts(list(electric.slides[9].shapes)[4]) == [
        "13.50 km",
        "  covered   \u00b7   SOC  ",
        "44.40 %",
        "   \u00b7   consumption  ",
        "1,234.5 Wh/km",
    ]

    hybrid_result = _with_values(
        hybrid_report.excel_result,
        {
            "engine_fuel_consumption_last": 8.75,
            "generator_power_1_max": 42.0,
            "engine_power_required_max": 55.0,
        },
    )
    hybrid_output = build_profile_powerpoint_report(hybrid_result, tmp_path / "hybrid_dynamic")
    hybrid_text = _visible_text(hybrid_output.presentation_path)
    assert "8.75 kg" in hybrid_text
    assert "42.00 kW" in hybrid_text
    assert "55.00 kW" in hybrid_text
    assert "RANGE EXTENDER INACTIVE IN THIS SIMULATION" not in hybrid_text


def test_profile_powerpoint_template_path_is_profile_configurable() -> None:
    electric_source = ELECTRIC_PROFILE.read_text(encoding="utf-8")
    hybrid_source = HYBRID_PROFILE.read_text(encoding="utf-8")
    engine_source = Path("src/vsm_postprocessing/profile_powerpoint_report_engine.py").read_text(encoding="utf-8")

    assert "powerpoint_template: reference_files/RoboSprayer_Electric_Report_FINAL.pptx" in electric_source
    assert "powerpoint_template: reference_files/RoboSprayer_Hybrid_Engineering_Report.pptx" in hybrid_source
    assert "fallback_mode: auxiliary_tyre_energy" not in electric_source
    assert "Hybrid_SP_Caiman_Sprayer_Report_FINAL (1).pptx" not in electric_source
    assert "Hybrid_SP_Caiman_Sprayer_Report_FINAL (1).pptx" not in hybrid_source
    for forbidden in ("3853", "12.00 km", "64.20 min", "2804.6", "17,418", "114.00 km", "290.28 min"):
        assert forbidden not in engine_source


def test_streamlit_profile_workflow_exposes_powerpoint_without_custom_analysis_regression() -> None:
    source = (PROJECT_ROOT / "src" / "vsm_postprocessing" / "ui_app.py").read_text(encoding="utf-8")

    assert "Custom Analysis" in source
    assert "generate_reporting_profile_engineering_report" in source
    assert "Download PowerPoint Engineering Report" in source
    assert "PowerPoint output is not available yet" not in source


def _slide_titles(path: Path) -> list[str]:
    prs = Presentation(path)
    titles = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = [getattr(shape, "text", "").strip() for shape in slide.shapes if getattr(shape, "text", "").strip()]
        titles.append(texts[1] if index == 1 else texts[0])
    return titles


def _reference_titles(path: Path) -> list[str]:
    prs = Presentation(path)
    slots = {
        "electric": {1: 3, 2: 2, 3: 2, 4: 2, 5: 2, 6: 2, 7: 2, 8: 2, 9: 2, 10: 2},
        "hybrid": {1: 2, 2: 0, 3: 0, 4: 0, 5: 0, 6: 0, 7: 0, 8: 0, 9: 0, 10: 1},
    }
    key = "electric" if "Electric" in path.name else "hybrid"
    return [list(slide.shapes)[slots[key][index]].text.strip() for index, slide in enumerate(prs.slides, start=1)]


def _first_run_font(shape):
    paragraph = shape.text_frame.paragraphs[0]
    return paragraph.runs[0].font if paragraph.runs else paragraph.font


def _run_texts(shape) -> list[str]:
    return _run_texts_from_paragraph(shape.text_frame.paragraphs[0])


def _run_texts_from_paragraph(paragraph) -> list[str]:
    return [run.text for run in paragraph.runs]


def _run_styles(shape) -> list[tuple[str | None, float | None, bool | None, str | None]]:
    styles = []
    for run in shape.text_frame.paragraphs[0].runs:
        font = run.font
        styles.append(
            (
                font.name,
                font.size.pt if font.size else None,
                font.bold,
                str(font.color.rgb) if getattr(font.color, "rgb", None) else None,
            )
        )
    return styles


def _find_text_shape(slide, text: str):
    return next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and text in shape.text)


def _shape_geometry(shape) -> tuple[int, int, int, int]:
    return shape.left, shape.top, shape.width, shape.height


def _picture_geometries(shapes) -> set[tuple[int, int, int, int]]:
    return {_shape_geometry(shape) for shape in shapes if shape.shape_type == 13}


def _media_hashes(path: Path) -> list[str]:
    with zipfile.ZipFile(path) as package:
        return sorted(
            hashlib.sha256(package.read(name)).hexdigest()
            for name in package.namelist()
            if name.startswith("ppt/media/") and not name.endswith("/")
        )


def _with_values(excel_result, replacements: dict[str, float]):
    result = copy.copy(excel_result)
    statistics_result = copy.copy(excel_result.statistics_result)
    statistics_result.statistics = [
        replace(item, value=replacements.get(item.definition.statistic_id, item.value))
        for item in statistics_result.statistics
    ]
    statistics_result.kpis = [
        replace(item, value=replacements.get(item.definition.kpi_id, item.value))
        for item in statistics_result.kpis
    ]
    result.statistics_result = statistics_result
    return result


def _visible_text(path: Path) -> str:
    prs = Presentation(path)
    return "\n".join(
        getattr(shape, "text", "")
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "text", "").strip()
    )
